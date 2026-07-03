"""Insert an official cover/divider/quote slide into an existing .pptx.

Written for sfnl-deck-edit. Sourced from engine/web/assets/chrome/manifest.json + its
PNGs — the same source of truth the HTML archetypes (engine/web/archetypes/) use for
chrome slides in the normal sfnl-deck build. Not a copy of raw slide XML from
engine/assets/sfnl-slides.pptx: that file's layouts/master/theme are foreign to an
arbitrary target deck, and reconciling them is exactly the fragile OOXML surgery this
plugin avoids. Painting the manifest's PNG as a full-bleed picture background and adding
plain text boxes at its slot coordinates reproduces the same visual result without any
of that risk.

Usage:
    python -m scripts.insert_chrome_slide target.pptx output.pptx <manifest-key> <position> \
        [--title "..."] [--subtitle "..."] [--body "..."]

<position> is the 0-based index the new slide lands at (0..slide_count, where
slide_count appends at the end). Text options map to whichever slot roles the chosen
key has; a key with no slots (e.g. cover-02) takes none. Omitted roles fall back to the
manifest's sample text.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

ENGINE = Path(__file__).resolve().parents[1]
MANIFEST_PATH = ENGINE / "web" / "assets" / "chrome" / "manifest.json"
CHROME_DIR = MANIFEST_PATH.parent

_ALIGNMENTS = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _load_manifest() -> dict[str, Any]:
    return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))


def _find_variant(manifest: dict[str, Any], key: str) -> dict[str, Any]:
    for variant in manifest["variants"]:
        if variant["key"] == key:
            return variant
    available = ", ".join(v["key"] for v in manifest["variants"])
    raise ValueError(f"Unknown chrome key {key!r}. Available: {available}")


def _strip_shapes(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def _move_slide(prs: Presentation, from_index: int, to_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    element = slide_id_list[from_index]
    slide_id_list.remove(element)
    slide_id_list.insert(to_index, element)


def _add_slot_textbox(slide, slot: dict[str, Any], text: str) -> None:
    box = slide.shapes.add_textbox(
        Pt(slot["left_pt"]), Pt(slot["top_pt"]), Pt(slot["width_pt"]), Pt(slot["height_pt"])
    )
    text_frame = box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    if slot.get("align") in _ALIGNMENTS:
        paragraph.alignment = _ALIGNMENTS[slot["align"]]
    run = paragraph.add_run()
    run.text = text
    run.font.name = slot["font"]
    run.font.size = Pt(slot["size_pt"])
    run.font.bold = bool(slot.get("bold"))
    run.font.color.rgb = RGBColor.from_string(slot["color"].lstrip("#").upper())


def insert_chrome_slide(
    pptx_path: str | Path,
    output_path: str | Path,
    key: str,
    position: int,
    texts: dict[str, str] | None = None,
) -> None:
    manifest = _load_manifest()
    variant = _find_variant(manifest, key)
    texts = texts or {}

    slot_roles = {slot["role"] for slot in variant["slots"]}
    unknown_roles = set(texts) - slot_roles
    if unknown_roles:
        raise ValueError(
            f"Chrome key {key!r} has no slot(s) for: {', '.join(sorted(unknown_roles))} "
            f"(available roles: {', '.join(sorted(slot_roles)) or 'none'})"
        )

    prs = Presentation(str(pptx_path))
    canvas_w_pt, canvas_h_pt = manifest["canvas_pt"]
    if (prs.slide_width.pt, prs.slide_height.pt) != (canvas_w_pt, canvas_h_pt):
        raise ValueError(
            f"Target slide size ({prs.slide_width.pt}x{prs.slide_height.pt}pt) doesn't "
            f"match the chrome manifest's canvas ({canvas_w_pt}x{canvas_h_pt}pt) — "
            "inserting would distort the official artwork."
        )

    slide_count = len(prs.slides)
    if not 0 <= position <= slide_count:
        raise ValueError(f"Position {position} out of range (0-{slide_count})")

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _strip_shapes(slide)
    slide.shapes.add_picture(
        str(CHROME_DIR / f"{key}.png"), 0, 0, prs.slide_width, prs.slide_height
    )
    for slot in variant["slots"]:
        _add_slot_textbox(slide, slot, texts.get(slot["role"], slot["sample"]))

    _move_slide(prs, slide_count, position)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("input", help="Target .pptx file")
    parser.add_argument("output", help="Output .pptx file")
    parser.add_argument("key", help="Chrome manifest key, e.g. divider-01")
    parser.add_argument("position", type=int, help="0-based insert position")
    parser.add_argument("--title")
    parser.add_argument("--subtitle")
    parser.add_argument("--body")
    args = parser.parse_args()

    texts = {
        role: value
        for role, value in {"title": args.title, "subtitle": args.subtitle, "body": args.body}.items()
        if value is not None
    }
    insert_chrome_slide(args.input, args.output, args.key, args.position, texts)
    print(f"Wrote {args.output}")


if __name__ == "__main__":
    main()
