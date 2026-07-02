"""Native vector icon library: pictograms drawn from pptx autoshape primitives.

No raster or SVG assets are bundled — every icon is composed at build time from
2-4 native autoshapes (oval, rectangle, triangle, ...) filled via `schemeClr` so
icons stay on-brand and scale cleanly at both small (icon-bubble) and large
(divider-panel) sizes, matching the reference deck's flat monochrome icon look
without adding an asset-embedding or licensing dependency.
"""
from __future__ import annotations

from math import cos, pi, sin

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from scripts.colors import set_run_font, set_run_scheme_color, set_scheme_fill, set_scheme_line

ICON_NAMES = (
    "target", "people", "growth", "idea", "house", "book", "calendar", "compass",
    "partnership", "check", "flag", "scale", "money", "clock", "gear",
)


def _rect(slide, x, y, w, h, fill, rotation=0.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_scheme_fill(shp, fill)
    shp.line.fill.background()
    if rotation:
        shp.rotation = rotation
    return shp


def _oval(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        set_scheme_fill(shp, fill)
    else:
        shp.fill.background()
    if line:
        set_scheme_line(shp, line, width_pt=1.5)
    else:
        shp.line.fill.background()
    return shp


def _triangle(slide, shape_type, x, y, w, h, fill, rotation=0.0):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_scheme_fill(shp, fill)
    shp.line.fill.background()
    if rotation:
        shp.rotation = rotation
    return shp


def _glyph_badge(slide, x, y, size, bg, glyph, glyph_color):
    _oval(slide, x, y, size, size, fill=bg)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(size), Inches(size))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    set_run_font(run, "Gotham Bold", size_pt=max(10, size * 26), bold=True)
    set_run_scheme_color(run, glyph_color)
    return box


def _target(slide, x, y, size, color, bg):
    _oval(slide, x, y, size, size, fill=color)
    _oval(slide, x + size * 0.19, y + size * 0.19, size * 0.62, size * 0.62, fill=bg)
    _oval(slide, x + size * 0.38, y + size * 0.38, size * 0.24, size * 0.24, fill=color)


def _people(slide, x, y, size, color, bg):
    _rect(slide, x + size * 0.06, y + size * 0.5, size * 0.4, size * 0.42, fill=color)
    _rect(slide, x + size * 0.54, y + size * 0.5, size * 0.4, size * 0.42, fill=color)
    _oval(slide, x + size * 0.1, y + size * 0.06, size * 0.32, size * 0.32, fill=color)
    _oval(slide, x + size * 0.58, y + size * 0.06, size * 0.32, size * 0.32, fill=color)


def _growth(slide, x, y, size, color, bg):
    bar_w = size * 0.22
    heights = (0.38, 0.62, 0.9)
    for i, h_frac in enumerate(heights):
        h = size * h_frac
        bx = x + size * 0.06 + i * (bar_w + size * 0.1)
        _rect(slide, bx, y + size - h, bar_w, h, fill=color)


def _idea(slide, x, y, size, color, bg):
    _oval(slide, x + size * 0.2, y, size * 0.6, size * 0.62, fill=color)
    _rect(slide, x + size * 0.36, y + size * 0.6, size * 0.28, size * 0.16, fill=color)
    _rect(slide, x + size * 0.4, y + size * 0.8, size * 0.2, size * 0.08, fill=color)


def _house(slide, x, y, size, color, bg):
    _triangle(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, size, size * 0.5, fill=color)
    _rect(slide, x + size * 0.16, y + size * 0.46, size * 0.68, size * 0.5, fill=color)


def _book(slide, x, y, size, color, bg):
    _rect(slide, x, y + size * 0.08, size * 0.9, size * 0.82, fill=color)
    _rect(slide, x + size * 0.14, y + size * 0.22, size * 0.62, size * 0.54, fill=bg)
    _rect(slide, x, y + size * 0.08, size * 0.1, size * 0.82, fill=color)


def _calendar(slide, x, y, size, color, bg):
    _rect(slide, x, y + size * 0.14, size, size * 0.82, fill=color)
    _rect(slide, x + size * 0.06, y + size * 0.32, size * 0.88, size * 0.56, fill=bg)
    _oval(slide, x + size * 0.14, y, size * 0.12, size * 0.24, fill=color)
    _oval(slide, x + size * 0.74, y, size * 0.12, size * 0.24, fill=color)


def _compass(slide, x, y, size, color, bg):
    _oval(slide, x, y, size, size, fill=bg, line=color)
    cx, cy = x + size / 2, y + size / 2
    needle_w, needle_h = size * 0.16, size * 0.38
    _triangle(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, cx - needle_w / 2, cy - needle_h, needle_w, needle_h,
              fill=color, rotation=0)
    _triangle(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, cx - needle_w / 2, cy, needle_w, needle_h,
              fill=color, rotation=180)


def _partnership(slide, x, y, size, color, bg):
    r = size * 0.34
    _rect(slide, x + size * 0.18, y + size * 0.42, size * 0.64, size * 0.16, fill=color)
    _oval(slide, x, y + (size - r) / 2, r, r, fill=color)
    _oval(slide, x + size - r, y + (size - r) / 2, r, r, fill=color)


def _check(slide, x, y, size, color, bg):
    _glyph_badge(slide, x, y, size, bg=color, glyph="✓", glyph_color=bg)


def _flag(slide, x, y, size, color, bg):
    _rect(slide, x + size * 0.08, y, size * 0.06, size, fill=color)
    _triangle(slide, MSO_SHAPE.RIGHT_TRIANGLE, x + size * 0.14, y + size * 0.05, size * 0.72, size * 0.42,
               fill=color)


def _scale(slide, x, y, size, color, bg):
    _rect(slide, x + size * 0.46, y + size * 0.1, size * 0.08, size * 0.7, fill=color)
    _rect(slide, x + size * 0.08, y + size * 0.1, size * 0.84, size * 0.07, fill=color)
    pan = size * 0.22
    _oval(slide, x + size * 0.02, y + size * 0.55, pan, pan * 0.7, fill=bg, line=color)
    _oval(slide, x + size * 0.76, y + size * 0.55, pan, pan * 0.7, fill=bg, line=color)
    _rect(slide, x + size * 0.02, y + size * 0.72, size * 0.96, size * 0.08, fill=color)


def _money(slide, x, y, size, color, bg):
    _glyph_badge(slide, x, y, size, bg=color, glyph="€", glyph_color=bg)


def _clock(slide, x, y, size, color, bg):
    _oval(slide, x, y, size, size, fill=bg, line=color)
    cx, cy = x + size / 2, y + size / 2
    _rect(slide, cx - size * 0.03, cy - size * 0.3, size * 0.06, size * 0.32, fill=color)
    _rect(slide, cx - size * 0.03, cy - size * 0.02, size * 0.24, size * 0.06, fill=color, rotation=0)
    _oval(slide, cx - size * 0.05, cy - size * 0.05, size * 0.1, size * 0.1, fill=color)


def _gear(slide, x, y, size, color, bg):
    cx, cy = x + size / 2, y + size / 2
    hub_r = size * 0.34
    _oval(slide, cx - hub_r / 2, cy - hub_r / 2, hub_r, hub_r, fill=color)
    tooth_w, tooth_h = size * 0.16, size * 0.2
    radius = size * 0.36
    for i in range(8):
        angle = 2 * pi * i / 8
        tx = cx + radius * cos(angle) - tooth_w / 2
        ty = cy + radius * sin(angle) - tooth_h / 2
        rotation = angle * 180 / pi
        _rect(slide, tx, ty, tooth_w, tooth_h, fill=color, rotation=rotation)


_BUILDERS = {
    "target": _target, "people": _people, "growth": _growth, "idea": _idea,
    "house": _house, "book": _book, "calendar": _calendar, "compass": _compass,
    "partnership": _partnership, "check": _check, "flag": _flag, "scale": _scale,
    "money": _money, "clock": _clock, "gear": _gear,
}


def draw_icon(slide, name: str, x: float, y: float, size: float, color: str = "orange", bg: str = "white"):
    """Draw a large monochrome pictogram in-place. Returns nothing; shapes are added to `slide`.

    `size` is the icon's square bounding box in inches. `color` is the ink `schemeClr`
    accent name; `bg` is the schemeClr used for negative space inside the icon (defaults
    to white — pass the panel's own background color when placing an icon on a colored
    full-bleed panel so the negative space reads correctly).
    """
    builder = _BUILDERS.get(str(name or "").strip().lower(), _target)
    builder(slide, x, y, size, color, bg)
