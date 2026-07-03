"""Extract the official SFNL cover/divider/quote slides as chrome archetypes.

Source of truth: ``engine/assets/sfnl-slides.pptx`` (a copy of the official
"SFNL Slides.pptx" sampler). Title slides and section dividers in generated
decks MUST use these designs — the pipeline never invents its own.

For every variant this script:

1. renders the source slide to a high-res PNG (PowerPoint COM, 2560x1440) with
   slide-number and date placeholders stripped, into
   ``engine/web/assets/chrome/<key>.png``;
2. reads the text-slot geometry and typography from the layout placeholders and
   writes ``engine/web/assets/chrome/manifest.json``;
3. generates one archetype HTML per variant in ``engine/web/archetypes/``,
   with the PNG as slide background and absolutely positioned text slots
   (class ``chrome-slot``) exactly where the sjabloon puts them.

Geometry mapping: the sjabloon canvas is 13.333in x 7.5in; the HTML canvas is
720pt x 405pt (10in x 5.625in), so 1 source inch = 54pt and font sizes scale
by 0.75.

Run from ``sfnl-pptx/engine``:

    python -m scripts.extract_chrome
"""
from __future__ import annotations

import copy
import json
import re
import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ENGINE = Path(__file__).resolve().parents[1]
SOURCE = ENGINE / "assets" / "sfnl-slides.pptx"
PALETTE = json.loads((ENGINE / "assets" / "palette.json").read_text(encoding="utf-8"))
CHROME_DIR = ENGINE / "web" / "assets" / "chrome"
ARCHETYPES = ENGINE / "web" / "archetypes"

PT_PER_SRC_IN = 54.0     # 720pt canvas / 13.333in source
FONT_SCALE = 0.75        # source pt -> deck pt

# bg/tx slots volgen de standaard clrMap van de master.
SCHEME_TO_SLOT = {"bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2"}

# key, 1-based source slide, kind, chrome mode advice, short catalog description
VARIANTS = [
    ("cover-01", 1, "cover", "none",
     "Fotocollage-kwadranten (oranje/royal/grapefruit + foto), wit logoblok in het midden; "
     "titel en subtitel wit rechtsonder."),
    ("cover-02", 2, "cover", "none",
     "Effen oranje vlak, wit logo rechtsonder. Geen tekstslots — opening of afsluiter."),
    ("cover-03", 3, "cover", "none",
     "Oranje cirkelcompositie, wit logo gecentreerd. Geen tekstslots — opening of afsluiter."),
    ("cover-04", 4, "cover", "none",
     "Oranje bogen-graphic bovenhelft, witte band onder; navy titel gecentreerd op de band."),
    ("quote-01", 5, "quote", "number",
     "Witte slide met grote foto en oranje overlayvlak, logo linksonder (in het ontwerp); "
     "titel en subtitel bovenaan."),
    ("divider-01", 6, "divider", "dark",
     "Foto links, royal vlak rechts; sectietitel boven met bodytekst (sectie-intro) eronder."),
    ("divider-02", 7, "divider", "dark",
     "Foto links (BMX-fietser voor stadspark), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-03", 8, "divider", "dark",
     "Foto links (straatbeeld, mannen voor snackbar), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-04", 9, "divider", "dark",
     "Foto links (hand met euro's op markt), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-05", 10, "divider", "dark",
     "Foto links (handen op laptop), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-06", 11, "divider", "dark",
     "Foto links (koffieoogst/landbouw), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-07", 12, "divider", "dark",
     "Foto links (demonstratiebord 'You paid for this'), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-08", 13, "divider", "dark",
     "Foto links (hardlopers met reflectie in plas), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-09", 14, "divider", "dark",
     "Foto links (brug over het IJ, Amsterdam), royal vlak rechts; titel + subtitel rechtsonder."),
    ("divider-10", 15, "divider", "dark",
     "Foto links (laptop met analytics-dashboard), royal vlak rechts; titel + subtitel rechtsonder."),
]

# placeholder idx -> (role, html tag, sample text) per layout family
SLOT_ROLES = {
    "1_Titelslide": {14: ("title", "h1", "DEKTITEL"),
                     13: ("subtitle", "p", "ONDERTITEL — KLANT — DATUM")},
    "7_Titelslide": {10: ("title", "h1", "DEKTITEL")},
    "Quote": {0: ("title", "h1", "TITEL"),
              1: ("subtitle", "p", "SUBTITEL")},
    "2_sectieslide_stijl1": {14: ("title", "h1", "SECTIETITEL"),
                             11: ("body", "p", "Sectie-intro in twee tot vier regels.")},
    "_sectie_default": {14: ("title", "h1", "SECTIETITEL"),
                        13: ("subtitle", "p", "SUBTITEL")},
}

FONT_NORMALIZE = {"Gotham Bold Regular": "Gotham Bold"}


def _resolve_color(token: str | None) -> str | None:
    """Map an srgbClr hex or schemeClr token to a #RRGGBB hex."""
    if token is None:
        return None
    if re.fullmatch(r"[0-9A-Fa-f]{6}", token):
        return f"#{token.upper()}"
    slot = SCHEME_TO_SLOT.get(token, token)
    entry = PALETTE["by_slot"].get(slot)
    return f"#{entry['hex']}" if entry else None


def _slot_style(layout_ph) -> dict:
    """First-paragraph typography of a layout placeholder (regex on its XML)."""
    xml = layout_ph._element.xml
    m = re.search(r'sz="(\d+)"', xml)
    size = round(int(m.group(1)) / 100 * FONT_SCALE, 2) if m else None
    m = re.search(r'typeface="([^"]+)"', xml)
    font = FONT_NORMALIZE.get(m.group(1), m.group(1)) if m else None
    m = re.search(r'(?:srgbClr|schemeClr) val="([^"]+)"', xml)
    color = _resolve_color(m.group(1) if m else None)
    m = re.search(r'algn="([^"]+)"', xml)
    align = {"ctr": "center", "r": "right"}.get(m.group(1)) if m else None
    bold = bool(re.search(r'b="1"', xml))
    return {"font": font, "size_pt": size, "color": color, "align": align, "bold": bold}


def _slots_for(slide) -> list[dict]:
    layout = slide.slide_layout
    roles = SLOT_ROLES.get(layout.name)
    if roles is None and layout.name.endswith("sectieslide_stijl1"):
        roles = SLOT_ROLES["_sectie_default"]
    if roles is None:
        return []
    lay_phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
    slots = []
    for sh in slide.placeholders:
        idx = sh.placeholder_format.idx
        if idx not in roles:
            continue
        role, tag, sample = roles[idx]
        geo_src = sh if sh.left is not None else lay_phs.get(idx)
        style = _slot_style(lay_phs[idx]) if idx in lay_phs else {}
        if not style.get("color"):
            # Layout zonder expliciete kleur (bv. Quote): erf de brand-defaults
            # voor donkere tekst op een witte slide.
            style["color"] = "#201B5C" if role == "title" else "#233348"
        left = max(0.0, Emu(geo_src.left).inches * PT_PER_SRC_IN)
        top = Emu(geo_src.top).inches * PT_PER_SRC_IN
        width = min(720.0 - left, Emu(geo_src.width).inches * PT_PER_SRC_IN)
        slots.append({
            "role": role, "tag": tag, "sample": sample,
            "left_pt": round(left, 1), "top_pt": round(top, 1),
            "width_pt": round(width, 1),
            "height_pt": round(Emu(geo_src.height).inches * PT_PER_SRC_IN, 1),
            **style,
        })
    order = {"title": 0, "subtitle": 1, "body": 2}
    slots.sort(key=lambda s: order.get(s["role"], 9))
    return slots


def _strip_number_and_date_placeholders(prs: Presentation) -> None:
    """Remove sldNum/dt placeholders and slide-number field shapes so the
    rendered PNGs carry no baked-in page number or date."""
    containers = list(prs.slide_masters)
    for master in prs.slide_masters:
        containers.extend(master.slide_layouts)
    for container in containers:
        spTree = container.shapes._spTree
        for sp in list(spTree):
            xml = getattr(sp, "xml", None)
            if xml is None:
                continue
            if ('type="sldNum"' in xml or 'type="dt"' in xml
                    or 'type="slidenum"' in xml or 'type="datetime"' in xml):
                spTree.remove(sp)


def _render_pngs(keys_by_slide: dict[int, str]) -> None:
    import pythoncom
    import win32com.client

    prs = Presentation(str(SOURCE))
    _strip_number_and_date_placeholders(prs)
    tmp = Path(tempfile.mkdtemp()) / "sfnl-slides-clean.pptx"
    prs.save(str(tmp))

    pythoncom.CoInitialize()
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(str(tmp), WithWindow=False)
            try:
                for slide_no, key in keys_by_slide.items():
                    dest = CHROME_DIR / f"{key}.png"
                    pres.Slides(slide_no).Export(str(dest), "PNG", 2560, 1440)
                    print("rendered", dest.name)
            finally:
                pres.Close()
        finally:
            app.Quit()
    finally:
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
    shutil.rmtree(tmp.parent, ignore_errors=True)


def _archetype_html(key: str, slots: list[dict]) -> str:
    lines = [
        "<!DOCTYPE html>",
        "<!-- Gegenereerd door scripts/extract_chrome.py uit engine/assets/sfnl-slides.pptx.",
        "     Niet met de hand herpositioneren: de slots volgen de officiële sjabloon-layout.",
        "     Kopieer engine/web/assets/chrome/ naar de workspace als slides/chrome/. -->",
        "<html>",
        '<head><meta charset="utf-8"><link rel="stylesheet" href="sfnl.css">',
        "<style>",
        f"  body {{ background-image: url('chrome/{key}.png'); background-size: cover; }}",
        "</style></head>",
        "<body>",
    ]
    for s in slots:
        style = [
            f"left: {s['left_pt']}pt", f"top: {s['top_pt']}pt", f"width: {s['width_pt']}pt",
            f"font-family: '{s['font']}', Arial, sans-serif",
            f"font-size: {s['size_pt']}pt",
            f"color: {s['color'] or '#FFFFFF'}",
        ]
        if s.get("align"):
            style.append(f"text-align: {s['align']}")
        if s.get("bold") and s["tag"] not in ("h1", "h2", "h3"):
            style.append("font-weight: bold")
        lines.append(
            f'  <{s["tag"]} class="chrome-slot" style="{"; ".join(style)};">{s["sample"]}</{s["tag"]}>'
        )
    lines += ["</body>", "</html>", ""]
    return "\n".join(lines)


def main() -> None:
    CHROME_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(SOURCE))

    manifest = {"source": "engine/assets/sfnl-slides.pptx", "canvas_pt": [720, 405],
                "variants": []}
    for key, slide_no, kind, chrome, desc in VARIANTS:
        slide = prs.slides[slide_no - 1]
        slots = _slots_for(slide)
        manifest["variants"].append({
            "key": key, "kind": kind, "source_slide": slide_no,
            "layout": slide.slide_layout.name, "chrome": chrome,
            "description": desc, "png": f"chrome/{key}.png", "slots": slots,
        })
        (ARCHETYPES / f"{key}.html").write_text(_archetype_html(key, slots), encoding="utf-8")
        print("wrote", f"archetypes/{key}.html", f"({len(slots)} slots)")

    (CHROME_DIR / "manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    print("wrote", "assets/chrome/manifest.json")

    _render_pngs({slide_no: key for key, slide_no, *_ in VARIANTS})


if __name__ == "__main__":
    main()
