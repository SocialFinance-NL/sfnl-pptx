"""Duplicate/reorder/delete slides within one .pptx by an index sequence.

Written for sfnl-deck-edit.

Usage:
    python -m scripts.rearrange input.pptx output.pptx 0,3,3,5

Produces output.pptx containing input.pptx's slides 0, 3 (twice — the repeat duplicates
it), and 5, in that order. Any original slide index not present in the sequence is
dropped. Indices are 0-based, into the *input* file's original slide order.
"""
from __future__ import annotations

import argparse
import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def _duplicate_slide(prs: Presentation, index: int):
    """Add a copy of prs.slides[index] at the end, returning the new slide."""
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    image_rels = {
        rid: rel for rid, rel in source.part.rels.items()
        if "image" in rel.reltype or "media" in rel.reltype
    }

    spTree = new_slide.shapes._spTree
    for shape in source.shapes:
        new_element = copy.deepcopy(shape._element)
        for blip in new_element.findall(f".//{qn('a:blip')}"):
            old_rid = blip.get(qn("r:embed"))
            rel = image_rels.get(old_rid)
            if rel is not None:
                blip.set(qn("r:embed"), new_slide.part.rels.get_or_add(rel.reltype, rel._target))
        spTree.append(new_element)

    return new_slide


def _delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id_element = slide_id_list[index]
    prs.part.drop_rel(slide_id_element.get(qn("r:id")))
    slide_id_list.remove(slide_id_element)


def rearrange(prs: Presentation, sequence: list[int]) -> None:
    total = len(prs.slides)
    for idx in sequence:
        if not 0 <= idx < total:
            raise ValueError(f"Slide index {idx} out of range (0-{total - 1})")

    # Resolve each position in the sequence to a concrete slide: the original for a
    # slide's first appearance, a fresh duplicate (appended) for every repeat.
    seen_once: set[int] = set()
    resolved_positions: list[int] = []  # positions into prs.slides *after* all duplication
    for original_idx in sequence:
        if original_idx in seen_once:
            _duplicate_slide(prs, original_idx)
            resolved_positions.append(len(prs.slides) - 1)  # duplicate is appended at the end
        else:
            seen_once.add(original_idx)
            resolved_positions.append(original_idx)

    # Drop original slides that no resolved position refers to (highest index first so
    # earlier indices stay valid while deleting).
    keep = set(resolved_positions)
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in keep:
            _delete_slide(prs, i)
            resolved_positions = [p - 1 if p > i else p for p in resolved_positions]

    # Rebuild the slide-ID list to match resolved_positions exactly, in one pass.
    slide_id_list = prs.slides._sldIdLst
    id_elements = list(slide_id_list)
    ordered = [id_elements[p] for p in resolved_positions]
    for el in id_elements:
        slide_id_list.remove(el)
    for el in ordered:
        slide_id_list.append(el)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("output", help="Output .pptx file")
    parser.add_argument("sequence", help="Comma-separated 0-based slide indices, e.g. 0,3,3,5")
    args = parser.parse_args()

    sequence = [int(x.strip()) for x in args.sequence.split(",")]
    prs = Presentation(args.input)
    rearrange(prs, sequence)

    Path(args.output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"Wrote {args.output} — {len(prs.slides)} slide(s)")


if __name__ == "__main__":
    main()
