"""Apply a text-replacement plan (from scripts.inventory's JSON shape) to a .pptx.

Written for sfnl-deck-edit. Every content shape found by scripts.inventory is cleared
unless the replacement plan gives it a "paragraphs" list — this is deliberate: an edit
plan must say explicitly what to keep, so stale text can't survive an edit by accident.

Does NOT auto-correct fonts/colors to SFNL brand defaults — only what the plan specifies
is applied. Run scripts.qa_text on the result to catch off-brand leftovers, and
scripts.render + visual inspection to catch any text overflow the edit introduced.

Usage:
    python -m scripts.replace input.pptx replacements.json output.pptx

Replacements JSON shape:
    {"slides": [{"index": 0, "shapes": [{"id": 0, "paragraphs": [{"text": "...", ...}]}]}]}

A shape entry with no "paragraphs" key (or omitted entirely) is cleared. Paragraph fields
mirror scripts.inventory's output: text, bullet, level, alignment, font_name, font_size_pt,
bold, italic, underline, color (hex RRGGBB), theme_color.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from scripts.inventory import iter_content_shapes

_ALIGNMENTS = {
    "LEFT": PP_ALIGN.LEFT,
    "CENTER": PP_ALIGN.CENTER,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
}
_BULLET_INDENT_PT = 18
# a:pPr children after the bu* group, per the DrawingML schema — a new bu* element must be
# inserted before these, not just appended, or PowerPoint may reject the paragraph.
_PPR_TAIL_TAGS = (qn("a:tabLst"), qn("a:defRPr"), qn("a:extLst"))


def _pPr_bullet_insert_index(pPr) -> int:
    for i, child in enumerate(pPr):
        if child.tag in _PPR_TAIL_TAGS:
            return i
    return len(pPr)


def _set_bullet(pPr, level: int) -> None:
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    pPr.set("marL", str(Pt(_BULLET_INDENT_PT * (level + 1))))
    pPr.set("indent", str(-Pt(_BULLET_INDENT_PT)))
    bu_char = pPr.makeelement(qn("a:buChar"), {"char": "•"})
    pPr.insert(_pPr_bullet_insert_index(pPr), bu_char)


def _clear_bullet(pPr) -> None:
    for tag in ("a:buChar", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    if pPr.find(qn("a:buNone")) is None:
        bu_none = pPr.makeelement(qn("a:buNone"), {})
        pPr.insert(_pPr_bullet_insert_index(pPr), bu_none)


def _apply_paragraph(paragraph, data: dict[str, Any]) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = data.get("text", "")

    pPr = paragraph._p.get_or_add_pPr()
    if data.get("bullet"):
        _set_bullet(pPr, data.get("level", 0))
    else:
        _clear_bullet(pPr)

    if "alignment" in data and data["alignment"] in _ALIGNMENTS:
        paragraph.alignment = _ALIGNMENTS[data["alignment"]]

    font = run.font
    if "font_name" in data:
        font.name = data["font_name"]
    if "font_size_pt" in data:
        font.size = Pt(data["font_size_pt"])
    if "bold" in data:
        font.bold = data["bold"]
    if "italic" in data:
        font.italic = data["italic"]
    if "underline" in data:
        font.underline = data["underline"]
    if "color" in data:
        font.color.rgb = RGBColor.from_string(data["color"].lstrip("#").upper())
    elif "theme_color" in data:
        try:
            font.color.theme_color = getattr(MSO_THEME_COLOR, data["theme_color"])
        except AttributeError as exc:
            raise ValueError(f"Unknown theme_color {data['theme_color']!r}") from exc


def _apply_shape(shape, paragraphs: list[dict[str, Any]] | None) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    if not paragraphs:
        return
    for i, para_data in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        _apply_paragraph(paragraph, para_data)


def apply_replacements(pptx_path: str | Path, plan: dict[str, Any], output_path: str | Path) -> None:
    prs = Presentation(str(pptx_path))

    plan_by_shape: dict[tuple[int, int], list[dict[str, Any]] | None] = {}
    for slide_plan in plan.get("slides", []):
        slide_index = slide_plan["index"]
        for shape_plan in slide_plan.get("shapes", []):
            plan_by_shape[(slide_index, shape_plan["id"])] = shape_plan.get("paragraphs")

    known = {(si, sid) for si, sid, *_ in iter_content_shapes(prs)}
    unknown = set(plan_by_shape) - known
    if unknown:
        lines = [f"  slide {si} shape {sid}" for si, sid in sorted(unknown)]
        raise ValueError("Replacement plan references shapes not found by inventory:\n" + "\n".join(lines))

    for slide_index, shape_id, shape, *_ in iter_content_shapes(prs):
        _apply_shape(shape, plan_by_shape.get((slide_index, shape_id)))

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("plan", help="Replacements .json file")
    parser.add_argument("output", help="Output .pptx file")
    args = parser.parse_args()

    plan = json.loads(Path(args.plan).read_text(encoding="utf-8"))
    apply_replacements(args.input, plan, args.output)
    print(f"Wrote {args.output}")


if __name__ == "__main__":
    main()
