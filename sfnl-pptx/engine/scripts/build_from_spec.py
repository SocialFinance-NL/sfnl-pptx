"""Deterministically build a .pptx from a validated deck-spec, on a clone of the
sjabloon. Template slides clone branded layouts and fill placeholders by idx;
custom slides draw on the sjabloon's `Titel, subtitel` layout (white canvas,
orange streepje, title/subtitle placeholders intact) with schemeClr colors.
Only full-bleed components (divider-block, closing-geometric) use `Leeg`.
Titles and subtitles are always rendered ALL CAPS."""
from __future__ import annotations

from math import cos, pi, sin
from pathlib import Path

from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from scripts.office.template import load_template_presentation
from scripts.extract_layouts import find_layout
from scripts.components import load_components
from scripts.spec import validate_spec, SpecError
from scripts.colors import (
    set_run_scheme_color, set_run_font, set_scheme_fill, set_scheme_line,
    set_scheme_fill_tint, set_scheme_fill_shade, set_fillformat_scheme, resolve_accent,
)
from scripts import icons

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5
FULL_BLEED_COMPONENTS = {"divider-block", "closing-geometric"}
CONTENT_LAYOUT = "Titel, subtitel"          # sjabloonlayout: wit vlak + oranje streepje + titel
CAPS_SLOTS = {"title", "subtitle"}          # titels en subtitels zijn altijd ALL CAPS
NO_CAPS_COMPONENTS = {"quote"}              # quote-'title' is lopende citaattekst, geen kop
VISUAL_COLORS = {"orange", "grapefruit", "royal", "sky", "emerald", "navy", "dark_slate", "white"}
ICON_GLYPHS = {
    "check": "✓",
    "gear": "⚙",
    "heart": "♡",
    "people": "○○",
    "target": "◎",
    "money": "€",
    "chart": "▥",
    "idea": "!",
}


def _placeholder_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx {idx} not found on slide")


def _set_placeholder_text(ph, value) -> None:
    tf = ph.text_frame
    if isinstance(value, list):
        tf.text = str(value[0]) if value else ""
        for item in value[1:]:
            p = tf.add_paragraph()
            p.text = str(item)
    else:
        tf.text = str(value)


def _color(value, fallback):
    return value if value in VISUAL_COLORS else fallback


def _caps(value):
    """Titels en subtitels staan huisstijlbreed in ALL CAPS."""
    if isinstance(value, str):
        return value.upper()
    if isinstance(value, list):
        return [_caps(v) for v in value]
    return value


def _glyph(value) -> str:
    key = str(value or "").strip()
    if not key:
        return "•"
    return ICON_GLYPHS.get(key.lower(), key[:2] if len(key) > 2 else key)


def _visual(comp, slide_spec, defaults: dict) -> dict:
    visual = dict(defaults)
    visual.update(comp.get("visual_defaults", {}))
    visual.update(slide_spec.get("visual") or {})
    return visual


def _emu(value) -> Emu:
    return Inches(float(value))


def _shape(slide, shape_type, x, y, w, h, fill="white", line=None):
    shp = slide.shapes.add_shape(shape_type, _emu(x), _emu(y), _emu(w), _emu(h))
    set_scheme_fill(shp, fill)
    if line:
        set_scheme_line(shp, line, width_pt=1.0)
    else:
        shp.line.fill.background()
    return shp


def _textbox(slide, text, x, y, w, h, font="Lato Light", size=15, color="navy",
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    set_run_font(run, font, size_pt=size, bold=bold)
    set_run_scheme_color(run, color)
    return box


def _icon_bubble(slide, icon, x, y, size, fill, text_color="white"):
    bubble = _shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill=fill)
    tf = bubble.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = _glyph(icon)
    set_run_font(run, "Gotham Bold", size_pt=max(12, float(size) * 22), bold=True)
    set_run_scheme_color(run, text_color)
    return bubble


def _line(slide, x1, y1, x2, y2, color="navy", width_pt=1.2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _emu(x1), _emu(y1), _emu(x2), _emu(y2))
    set_scheme_line(connector, color, width_pt=width_pt)
    return connector


def _tint_shape(slide, shape_type, x, y, w, h, accent, tint_pct=85, line=None):
    shp = slide.shapes.add_shape(shape_type, _emu(x), _emu(y), _emu(w), _emu(h))
    set_scheme_fill_tint(shp, accent, tint_pct)
    if line:
        set_scheme_line(shp, line, width_pt=1.0)
    else:
        shp.line.fill.background()
    return shp




def _build_template_slide(prs, comp, slide_spec):
    layout = find_layout(prs, comp["source_layout"])
    slide = prs.slides.add_slide(layout)
    fill = slide_spec.get("content_schema_fill", {})
    for slot_name, slot_def in comp["slots"].items():
        if slot_name in fill:
            value = fill[slot_name]
            if slot_name in CAPS_SLOTS and comp["id"] not in NO_CAPS_COMPONENTS:
                value = _caps(value)
            ph = _placeholder_by_idx(slide, slot_def["placeholder_idx"])
            _set_placeholder_text(ph, value)
    return slide


def _fill_content_frame(slide, fill) -> None:
    """Fill the sjabloon's own title/subtitle placeholders on the CONTENT_LAYOUT.
    The white canvas, orange streepje, and title styling come from the layout itself
    and stay intact; we only put text in the placeholders (always ALL CAPS)."""
    title_ph = _placeholder_by_idx(slide, 0)
    _set_placeholder_text(title_ph, _caps(str(fill.get("title", ""))))
    subtitle_ph = _placeholder_by_idx(slide, 1)
    subtitle = fill.get("subtitle")
    if subtitle:
        _set_placeholder_text(subtitle_ph, _caps(str(subtitle)))
    else:
        subtitle_ph._element.getparent().remove(subtitle_ph._element)


def _render_content_cards(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    cards = fill.get("cards", [])[:3]
    x0 = float(visual.get("x", 0.85)); y = float(visual.get("y", 2.15))
    w = float(visual.get("card_width", 3.35)); h = float(visual.get("card_height", 2.55))
    gap = float(visual.get("gap", 0.38)); icon_size = float(visual.get("icon_size", 0.5))
    header_h = min(0.7, h * 0.32)
    pastel = visual.get("variant") == "pastel-tint"
    for i, card in enumerate(cards):
        x = x0 + i * (w + gap)
        color = _color(card.get("color"), accent)
        if pastel:
            _tint_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, color, tint_pct=85, line=color)
        else:
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill="white", line=color)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, header_h, fill=color)
        _icon_bubble(slide, card.get("icon"), x + 0.18, y + 0.15, icon_size, fill="white", text_color=color)
        _textbox(slide, card.get("heading", ""), x + 0.25 + icon_size, y + 0.13, w - 0.45 - icon_size,
                 header_h - 0.08, font="Gotham Bold", size=15, color="white", bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        _textbox(slide, card.get("body", ""), x + 0.28, y + header_h + 0.25, w - 0.56, h - header_h - 0.35,
                 font="Lato Light", size=14, color="navy")


def _render_kpi_trio(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    kpis = fill.get("kpis", [])[:4]
    x0 = float(visual.get("x", 0.65)); y = float(visual.get("y", 2.25))
    w = float(visual.get("card_width", 3.0)); h = float(visual.get("card_height", 1.65))
    gap = float(visual.get("gap", 0.3))
    for i, kpi in enumerate(kpis):
        x = x0 + i * (w + gap)
        color = _color(kpi.get("color"), accent)
        value = str(kpi.get("value", ""))
        value_size = 27
        if len(value) > 14:
            value_size = 18
        elif len(value) > 10:
            value_size = 21
        elif len(value) > 7:
            value_size = 24
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill="white", line=color)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.12, fill=color)
        _textbox(slide, value, x + 0.12, y + 0.36, w - 0.24, 0.42,
                 font="Gotham Bold", size=value_size, color=color, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, kpi.get("label", ""), x + 0.2, y + 1.05, w - 0.4, 0.38,
                 font="Lato Light", size=12.5, color="navy", align=PP_ALIGN.CENTER)
    progress = visual.get("progress")
    if progress is not None:
        pct = max(0.0, min(1.0, float(progress)))
        bar_x = x0
        bar_y = y + h + 0.58
        bar_w = max(1.0, len(kpis) * w + max(0, len(kpis) - 1) * gap)
        _shape(slide, MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_w, 0.28, fill="sky")
        _shape(slide, MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_w * pct, 0.28, fill=accent)
        if visual.get("progress_label"):
            _textbox(slide, visual["progress_label"], bar_x, bar_y + 0.42, bar_w, 0.45,
                     font="Gotham Bold", size=16, color="navy", bold=True, align=PP_ALIGN.CENTER)


def _render_chart_static(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    series = fill.get("series", [])
    chart_x = float(visual.get("x", 1.0)); chart_y = float(visual.get("y", 2.0))
    chart_w = float(visual.get("width", 10.9)); chart_h = float(visual.get("height", 3.55))
    max_v = max((s.get("value", 0) for s in series), default=1) or 1
    if visual.get("show_grid", True):
        for i in range(5):
            gy = chart_y + chart_h * i / 4
            _shape(slide, MSO_SHAPE.RECTANGLE, chart_x, gy, chart_w, 0.01, fill="sky")
    n = max(len(series), 1)
    gap = 0.35
    bar_w = max(0.25, (chart_w - gap * (n - 1)) / n * 0.58)
    step = chart_w / n
    for i, s in enumerate(series):
        h = chart_h * (float(s.get("value", 0)) / max_v)
        x = chart_x + i * step + (step - bar_w) / 2
        color = _color(s.get("color"), accent)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, chart_y + chart_h - h, bar_w, h, fill=color)
        _textbox(slide, s.get("label", ""), x - 0.15, chart_y + chart_h + 0.18, bar_w + 0.3, 0.32,
                 font="Lato Light", size=11, color="navy", align=PP_ALIGN.CENTER)


def _render_process_timeline(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    steps = fill.get("steps", [])[:5]
    x0 = float(visual.get("x", 0.9)); y = float(visual.get("y", 2.15))
    w = float(visual.get("step_width", 2.45)); h = float(visual.get("step_height", 1.05))
    gap = float(visual.get("gap", 0.22)); icon_size = float(visual.get("icon_size", 0.42))
    shape_type = MSO_SHAPE.CHEVRON if visual.get("variant") == "chevron" else MSO_SHAPE.RIGHT_ARROW
    for i, step in enumerate(steps):
        x = x0 + i * (w + gap)
        color = _color(step.get("color"), accent)
        arrow = _shape(slide, shape_type, x, y, w, h, fill=color)
        arrow.text_frame.margin_left = Inches(0.12)
        _icon_bubble(slide, step.get("icon", i + 1), x + 0.18, y + 0.28, icon_size, fill="white", text_color=color)
        _textbox(slide, step.get("label", ""), x + 0.28 + icon_size, y + 0.17, w - 0.62 - icon_size, 0.34,
                 font="Gotham Bold", size=10.5, color="white", bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, step.get("detail", ""), x + 0.28 + icon_size, y + 0.54, w - 0.62 - icon_size, 0.34,
                 font="Lato Light", size=8.5, color="white", align=PP_ALIGN.CENTER)


def _node_position(index, visual):
    columns = max(1, int(visual.get("columns", 4)))
    w = float(visual.get("box_width", 2.35)); h = float(visual.get("box_height", 1.2))
    gap = float(visual.get("gap", 0.32)); row_gap = float(visual.get("row_gap", 0.38))
    x = float(visual.get("x", 0.8)) + (index % columns) * (w + gap)
    y = float(visual.get("y", 2.0)) + (index // columns) * (h + row_gap)
    return x, y, w, h


def _render_schema_grid(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    nodes = fill.get("nodes", [])[:12]
    for connection in fill.get("connections", []):
        try:
            source = int(connection.get("from"))
            target = int(connection.get("to"))
        except (TypeError, ValueError):
            continue
        if source >= len(nodes) or target >= len(nodes) or source < 0 or target < 0:
            continue
        sx, sy, sw, sh = _node_position(source, visual)
        tx, ty, tw, th = _node_position(target, visual)
        _line(slide, sx + sw, sy + sh / 2, tx, ty + th / 2, color=accent, width_pt=1.4)
    icon_size = float(visual.get("icon_size", 0.34))
    for i, node in enumerate(nodes):
        x, y, w, h = _node_position(i, visual)
        color = _color(node.get("color"), accent)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill="white", line=color)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.36, fill=color)
        _icon_bubble(slide, node.get("icon"), x + 0.12, y + 0.07, icon_size, fill="white", text_color=color)
        _textbox(slide, node.get("heading", ""), x + 0.2 + icon_size, y + 0.07, w - 0.32 - icon_size, 0.26,
                 font="Gotham Bold", size=10.5, color="white", bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, node.get("body", ""), x + 0.18, y + 0.52, w - 0.36, h - 0.62,
                 font="Lato Light", size=10, color="navy")


def _render_image_icon_trio(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    items = fill.get("items", [])[:3]
    x0 = float(visual.get("x", 0.75)); y = float(visual.get("y", 1.95))
    w = float(visual.get("column_width", 3.45)); img_h = float(visual.get("image_height", 1.6))
    gap = float(visual.get("gap", 0.42)); icon_size = float(visual.get("icon_size", 0.52))
    for i, item in enumerate(items):
        x = x0 + i * (w + gap)
        color = _color(item.get("color"), accent)
        image = item.get("image")
        if image and Path(image).exists():
            slide.shapes.add_picture(str(Path(image)), _emu(x), _emu(y + 0.88), width=_emu(w), height=_emu(img_h))
        else:
            _shape(slide, MSO_SHAPE.RECTANGLE, x, y + 0.88, w, img_h, fill=color)
        _icon_bubble(slide, item.get("icon"), x + 0.2, y, icon_size, fill=color)
        _textbox(slide, item.get("heading", ""), x + 0.9, y + 0.05, w - 0.9, 0.45,
                 font="Gotham Bold", size=16, color="navy", bold=True)
        _textbox(slide, item.get("body", ""), x, y + 0.56, w, 0.32, font="Lato Light", size=12.5, color="navy")


def _render_matrix_2x2(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    x = float(visual.get("x", 0.95)); y = float(visual.get("y", 1.95))
    w = float(visual.get("width", 6.8)); h = float(visual.get("height", 4.2))
    label_w = float(visual.get("label_width", 1.0))
    cell_w = w / 2
    cell_h = h / 2
    quadrants = fill.get("quadrants", [])[:4]
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, quadrant in enumerate(quadrants):
        col, row = positions[i]
        qx = x + col * cell_w
        qy = y + row * cell_h
        color = _color(quadrant.get("color"), accent)
        _shape(slide, MSO_SHAPE.RECTANGLE, qx, qy, cell_w, cell_h, fill="white", line=color)
        _shape(slide, MSO_SHAPE.RECTANGLE, qx, qy, cell_w, 0.22, fill=color)
        _textbox(slide, quadrant.get("heading", ""), qx + 0.2, qy + 0.45, cell_w - 0.4, 0.35,
                 font="Gotham Bold", size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, quadrant.get("body", ""), qx + 0.25, qy + 0.95, cell_w - 0.5, 0.55,
                 font="Lato Light", size=12.5, color="navy", align=PP_ALIGN.CENTER)
    _line(slide, x + cell_w, y, x + cell_w, y + h, color="navy", width_pt=1.0)
    _line(slide, x, y + cell_h, x + w, y + cell_h, color="navy", width_pt=1.0)
    _textbox(slide, fill.get("y_axis", ""), x - label_w, y + h / 2 - 0.18, label_w - 0.12, 0.36,
             font="Gotham Bold", size=12, color="navy", bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, fill.get("x_axis", ""), x + w / 2 - 0.55, y + h + 0.2, 1.1, 0.35,
             font="Gotham Bold", size=12, color="navy", bold=True, align=PP_ALIGN.CENTER)


def _render_layer_stack(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    layers = fill.get("layers", [])[:7]
    x0 = float(visual.get("x", 1.0)); y0 = float(visual.get("y", 2.0))
    w0 = float(visual.get("width", 8.8)); h = float(visual.get("layer_height", 0.58))
    gap = float(visual.get("gap", 0.08)); indent = float(visual.get("indent", 0.36))
    for i, layer in enumerate(layers):
        x = x0 + i * indent
        w = max(2.0, w0 - i * indent * 2)
        y = y0 + i * (h + gap)
        color = _color(layer.get("color"), accent)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=color)
        _textbox(slide, layer.get("label", ""), x + 0.25, y + 0.14, 1.35, 0.28,
                 font="Gotham Bold", size=10.5, color="white", bold=True)
        _textbox(slide, layer.get("detail", ""), x + 1.65, y + 0.14, w - 1.9, 0.28,
                 font="Lato Light", size=10.5, color="white")


def _render_cycle_diagram(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    steps = fill.get("steps", [])[:6]
    cx = float(visual.get("center_x", 6.65)); cy = float(visual.get("center_y", 3.9))
    radius = float(visual.get("radius", 1.65)); node_size = float(visual.get("node_size", 0.92))
    centers = []
    for i, _step in enumerate(steps):
        angle = -pi / 2 + (2 * pi * i / max(1, len(steps)))
        centers.append((cx + radius * cos(angle), cy + radius * sin(angle)))
    for i in range(len(centers)):
        x1, y1 = centers[i]
        x2, y2 = centers[(i + 1) % len(centers)]
        _line(slide, x1, y1, x2, y2, color=accent, width_pt=1.2)
    _shape(slide, MSO_SHAPE.OVAL, cx - 0.72, cy - 0.72, 1.44, 1.44, fill="white", line=accent)
    _textbox(slide, fill.get("center", ""), cx - 0.58, cy - 0.16, 1.16, 0.32,
             font="Gotham Bold", size=13, color=accent, bold=True, align=PP_ALIGN.CENTER)
    for i, step in enumerate(steps):
        nx, ny = centers[i]
        color = _color(step.get("color"), accent)
        _shape(slide, MSO_SHAPE.OVAL, nx - node_size / 2, ny - node_size / 2, node_size, node_size, fill=color)
        _textbox(slide, _glyph(step.get("icon", i + 1)), nx - 0.18, ny - 0.17, 0.36, 0.28,
                 font="Gotham Bold", size=12, color="white", bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, step.get("label", ""), nx - 0.75, ny + node_size / 2 + 0.08, 1.5, 0.32,
                 font="Gotham Bold", size=11, color="navy", bold=True, align=PP_ALIGN.CENTER)


def _render_scenario_cards(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    scenarios = fill.get("scenarios", [])[:3]
    x0 = float(visual.get("x", 0.75)); y = float(visual.get("y", 1.85))
    w = float(visual.get("card_width", 3.55)); h = float(visual.get("card_height", 4.25))
    gap = float(visual.get("gap", 0.35))
    for i, scenario in enumerate(scenarios):
        x = x0 + i * (w + gap)
        color = _color(scenario.get("color"), accent)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill="white", line=color)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.56, fill=color)
        _textbox(slide, scenario.get("heading", ""), x + 0.22, y + 0.13, w - 0.44, 0.28,
                 font="Gotham Bold", size=12, color="white", bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, scenario.get("subtitle", ""), x + 0.24, y + 0.78, w - 0.48, 0.34,
                 font="Gotham Bold", size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        _shape(slide, MSO_SHAPE.RECTANGLE, x + 0.25, y + 1.25, 0.44, 0.045, fill=color)
        _textbox(slide, "Kenmerken", x + 0.25, y + 1.35, w - 0.5, 0.24,
                 font="Gotham Bold", size=10.5, color="navy", bold=True)
        for j, item in enumerate(scenario.get("pros", [])[:3]):
            _textbox(slide, "- " + str(item), x + 0.32, y + 1.68 + j * 0.3, w - 0.64, 0.22,
                     font="Lato Light", size=9.6, color="navy")
        _shape(slide, MSO_SHAPE.RECTANGLE, x + 0.25, y + 2.65, 0.44, 0.045, fill=color)
        _textbox(slide, "Aandachtspunten", x + 0.25, y + 2.75, w - 0.5, 0.24,
                 font="Gotham Bold", size=10.5, color=color, bold=True)
        for j, item in enumerate(scenario.get("cons", [])[:3]):
            _textbox(slide, "- " + str(item), x + 0.32, y + 3.08 + j * 0.3, w - 0.64, 0.22,
                     font="Lato Light", size=9.6, color="navy")


def _render_assessment_table(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    rows = fill.get("rows", [])[:8]
    columns = fill.get("columns", ["Criterium", "Toelichting", "Score"])
    x = float(visual.get("x", 0.75)); y = float(visual.get("y", 1.8))
    w = float(visual.get("width", 10.9)); row_h = float(visual.get("row_height", 0.62))
    col_ws = [w * 0.28, w * 0.52, w * 0.20]
    header_h = 0.48
    cur_x = x
    for i, header in enumerate((columns + ["", "", ""])[:3]):
        _shape(slide, MSO_SHAPE.RECTANGLE, cur_x, y, col_ws[i], header_h, fill="navy")
        _textbox(slide, header, cur_x + 0.12, y + 0.12, col_ws[i] - 0.24, 0.22,
                 font="Gotham Bold", size=10, color="white", bold=True)
        cur_x += col_ws[i]
    for r, row in enumerate(rows):
        row_y = y + header_h + r * row_h
        color = _color(row.get("color"), accent)
        cur_x = x
        values = [row.get("criterion", ""), row.get("description", ""), row.get("score", "")]
        for c, value in enumerate(values):
            fill_color = color if c == 2 else "white"
            text_color = "white" if c == 2 else "navy"
            _shape(slide, MSO_SHAPE.RECTANGLE, cur_x, row_y, col_ws[c], row_h, fill=fill_color, line=color)
            _textbox(slide, value, cur_x + 0.12, row_y + 0.17, col_ws[c] - 0.24, 0.25,
                     font="Gotham Bold" if c in (0, 2) else "Lato Light",
                     size=9.6, color=text_color, bold=(c in (0, 2)),
                     align=PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT)
            cur_x += col_ws[c]


def _flow_endpoint(value, center_pos, node_positions):
    if value == "center":
        return center_pos
    try:
        return node_positions[int(value)]
    except (TypeError, ValueError, IndexError):
        return center_pos


def _render_mechanism_diagram(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    cx = float(visual.get("center_x", 6.4)); cy = float(visual.get("center_y", 3.75))
    node_w = float(visual.get("node_width", 2.0)); node_h = float(visual.get("node_height", 0.82))
    rx = float(visual.get("radius_x", 3.6)); ry = float(visual.get("radius_y", 1.75))
    nodes = fill.get("nodes", [])[:8]
    node_positions = []
    for i, _node in enumerate(nodes):
        angle = -pi / 2 + (2 * pi * i / max(1, len(nodes)))
        node_positions.append((cx + rx * cos(angle), cy + ry * sin(angle)))
    center_pos = (cx, cy)
    for flow in fill.get("flows", []):
        x1, y1 = _flow_endpoint(flow.get("from"), center_pos, node_positions)
        x2, y2 = _flow_endpoint(flow.get("to"), center_pos, node_positions)
        _line(slide, x1, y1, x2, y2, color=accent, width_pt=1.1)
        if flow.get("label"):
            label_x = (x1 + x2) / 2 - 0.45
            label_y = (y1 + y2) / 2 - 0.17
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, label_x - 0.04, label_y - 0.02, 0.98, 0.34,
                   fill="white", line=accent)
            _textbox(slide, flow["label"], label_x, label_y, 0.9, 0.3,
                     font="Gotham Bold", size=8.5, color=accent, bold=True, align=PP_ALIGN.CENTER)
    center = fill.get("center", {})
    center_color = _color(center.get("color"), accent)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - node_w / 2, cy - node_h / 2, node_w, node_h,
           fill=center_color)
    _textbox(slide, center.get("heading", ""), cx - node_w / 2 + 0.15, cy - 0.24, node_w - 0.3, 0.22,
             font="Gotham Bold", size=10.5, color="white", bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, center.get("body", ""), cx - node_w / 2 + 0.15, cy + 0.04, node_w - 0.3, 0.22,
             font="Lato Light", size=8.5, color="white", align=PP_ALIGN.CENTER)
    for i, node in enumerate(nodes):
        nx, ny = node_positions[i]
        color = _color(node.get("color"), accent)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, nx - node_w / 2, ny - node_h / 2, node_w, node_h,
               fill="white", line=color)
        _icon_bubble(slide, node.get("icon"), nx - node_w / 2 + 0.12, ny - 0.18, 0.36, fill=color)
        _textbox(slide, node.get("heading", ""), nx - node_w / 2 + 0.56, ny - 0.26, node_w - 0.68, 0.22,
                 font="Gotham Bold", size=8.8, color=color, bold=True)
        _textbox(slide, node.get("body", ""), nx - node_w / 2 + 0.56, ny + 0.03, node_w - 0.68, 0.22,
                 font="Lato Light", size=8.0, color="navy")


def _render_divider_block(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    panel_w = float(visual.get("panel_width", 5.3))
    icon_size = float(visual.get("icon_size", 2.3))
    side = visual.get("panel_side", "left")
    panel_x = 0.0 if side == "left" else SLIDE_W_IN - panel_w
    _shape(slide, MSO_SHAPE.RECTANGLE, panel_x, 0, panel_w, SLIDE_H_IN, fill=accent)
    icons.draw_icon(slide, fill.get("icon"), panel_x + (panel_w - icon_size) / 2,
                     (SLIDE_H_IN - icon_size) / 2, icon_size, color="white", bg=accent)
    text_x = panel_w + 0.75 if side == "left" else 0.75
    text_w = SLIDE_W_IN - panel_w - 1.5
    _textbox(slide, _caps(fill.get("title", "")), text_x, 2.85, text_w, 1.3,
             font="Gotham Bold", size=30, color="navy", bold=True)
    _shape(slide, MSO_SHAPE.RECTANGLE, text_x, 3.95, 0.5, 0.06, fill=accent)
    if fill.get("subtitle"):
        _textbox(slide, _caps(fill["subtitle"]), text_x, 4.15, text_w, 1.1, font="Lato Light", size=14, color="navy")


def _render_stat_banner(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    banners = fill.get("banners", [])[:3]
    x = float(visual.get("x", 0.85)); y0 = float(visual.get("y", 2.1))
    w = float(visual.get("width", 11.6)); h = float(visual.get("banner_height", 1.15))
    gap = float(visual.get("gap", 0.22))
    for i, banner in enumerate(banners):
        y = y0 + i * (h + gap)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill="navy")
        if banner.get("context"):
            _textbox(slide, banner["context"], x + 0.3, y, w * 0.55, h,
                     font="Lato Light", size=12, color="white", anchor=MSO_ANCHOR.MIDDLE)
        color = _color(banner.get("color"), accent)
        value_x = x + w * 0.55
        _textbox(slide, str(banner.get("value", "")), value_x, y + 0.1, w * 0.4, h * 0.62,
                 font="Gotham Bold", size=30, color=color, bold=True, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.BOTTOM)
        _textbox(slide, banner.get("label", ""), value_x, y + h * 0.64, w * 0.4, h * 0.32,
                 font="Lato Light", size=11, color="white", align=PP_ALIGN.RIGHT)


def _render_swimlane_columns(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    columns = fill.get("columns", [])[:6]
    x0 = float(visual.get("x", 0.7)); y = float(visual.get("y", 1.95))
    width = float(visual.get("width", 11.9)); height = float(visual.get("height", 4.85))
    header_h = float(visual.get("header_height", 0.52)); gap = float(visual.get("gap", 0.14))
    n = max(len(columns), 1)
    col_w = (width - gap * (n - 1)) / n
    for i, col in enumerate(columns):
        cx = x0 + i * (col_w + gap)
        color = _color(col.get("color"), accent)
        _shape(slide, MSO_SHAPE.RECTANGLE, cx, y, col_w, header_h, fill=color)
        _textbox(slide, col.get("heading", ""), cx + 0.08, y, col_w - 0.16, header_h,
                 font="Gotham Bold", size=11, color="white", bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        _tint_shape(slide, MSO_SHAPE.RECTANGLE, cx, y + header_h, col_w, height - header_h, color, tint_pct=85)
        items = col.get("items", [])[:8]
        if items:
            item_h = max(0.28, (height - header_h - 0.2) / len(items))
            item_y = y + header_h + 0.14
            for item in items:
                _textbox(slide, "• " + str(item), cx + 0.12, item_y, col_w - 0.22, item_h,
                         font="Lato Light", size=9.5, color="navy")
                item_y += item_h


def _render_closing_geometric(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=accent)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, _emu(9.3), _emu(-1.3), _emu(4.8), _emu(4.8))
    set_scheme_fill_tint(circle, accent, 35)
    circle.line.fill.background()
    tri1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _emu(0.4), _emu(5.4), _emu(2.6), _emu(2.6))
    set_scheme_fill_shade(tri1, accent, 20)
    tri1.line.fill.background()
    tri1.rotation = 15
    tri2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _emu(1.7), _emu(6.3), _emu(1.5), _emu(1.5))
    set_scheme_fill(tri2, "white")
    tri2.line.fill.background()
    tri2.rotation = -10
    _textbox(slide, _caps(fill.get("title", "")), 1.2, 3.0, 10.9, 1.0,
             font="Gotham Bold", size=32, color="white", bold=True, align=PP_ALIGN.CENTER)
    if fill.get("subtitle"):
        _textbox(slide, _caps(fill["subtitle"]), 1.2, 4.0, 10.9, 0.6,
                 font="Lato Light", size=16, color="white", align=PP_ALIGN.CENTER)


_FREEFORM_SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded-rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "chevron": MSO_SHAPE.CHEVRON,
    "donut": MSO_SHAPE.DONUT,
    "pie": MSO_SHAPE.PIE,
}
_ARROW_SHAPES = {
    "right": MSO_SHAPE.RIGHT_ARROW,
    "left": MSO_SHAPE.LEFT_ARROW,
    "up": MSO_SHAPE.UP_ARROW,
    "down": MSO_SHAPE.DOWN_ARROW,
}
_FREEFORM_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_FREEFORM_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _connector_arrowheads(connector, arrow: str) -> None:
    """Add triangle arrowhead(s) to a connector via ln XML (schemeClr already set)."""
    from pptx.oxml.ns import qn
    ln = connector.line._get_or_add_ln()
    if arrow in ("end", "both"):
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
        ln.append(tail)
    if arrow == "both":
        head = ln.makeelement(qn("a:headEnd"), {"type": "triangle"})
        ln.append(head)


def _freeform_connector(shapes_host, prim, accent):
    style = MSO_CONNECTOR.ELBOW if prim.get("style") == "elbow" else MSO_CONNECTOR.STRAIGHT
    connector = shapes_host.shapes.add_connector(
        style, _emu(prim["x1"]), _emu(prim["y1"]), _emu(prim["x2"]), _emu(prim["y2"]))
    set_scheme_line(connector, prim.get("color") or accent, width_pt=float(prim.get("width_pt", 1.4)))
    _connector_arrowheads(connector, prim.get("arrow", "end"))
    return connector


def _freeform_textbox(shapes_host, prim, accent):
    align = _FREEFORM_ALIGN.get(prim.get("align", "left"), PP_ALIGN.LEFT)
    anchor = _FREEFORM_ANCHOR.get(prim.get("anchor", "top"), MSO_ANCHOR.TOP)
    font = prim.get("font", "Lato Light")
    size = float(prim.get("size", 14))
    color = prim.get("color", "navy")
    bold = bool(prim.get("bold", False))
    bullets = prim.get("bullets")
    if bullets:
        box = shapes_host.shapes.add_textbox(_emu(prim["x"]), _emu(prim["y"]), _emu(prim["w"]), _emu(prim["h"]))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = "• " + str(item)
            set_run_font(run, font, size_pt=size, bold=bold)
            set_run_scheme_color(run, color)
        return box
    return _textbox(shapes_host, prim.get("text", ""), prim["x"], prim["y"], prim["w"], prim["h"],
                    font=font, size=size, color=color, bold=bold, align=align, anchor=anchor)


def _freeform_table(shapes_host, prim, accent):
    rows = prim.get("rows", [])
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    header = prim.get("header", True)
    row_h = float(prim.get("row_height", 0.42))
    x, y, w = float(prim["x"]), float(prim["y"]), float(prim["w"])
    frame = shapes_host.shapes.add_table(n_rows, n_cols, _emu(x), _emu(y), _emu(w), _emu(n_rows * row_h))
    table = frame.table
    table.first_row = bool(header)
    table.horz_banding = False
    col_widths = prim.get("col_widths")
    if isinstance(col_widths, list) and len(col_widths) == n_cols:
        total = sum(float(c) for c in col_widths) or 1.0
        for c, colw in enumerate(col_widths):
            table.columns[c].width = _emu(w * float(colw) / total)
    accent_color = prim.get("color") or accent
    for r, row in enumerate(rows):
        is_header = header and r == 0
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            set_fillformat_scheme(cell.fill, "navy" if is_header else "white")
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(value)
            if is_header:
                set_run_font(run, "Gotham Bold", size_pt=float(prim.get("header_size", 10.5)), bold=True)
                set_run_scheme_color(run, "white")
            else:
                first_col_bold = bool(prim.get("first_col_bold", False)) and c == 0
                set_run_font(run, "Gotham Bold" if first_col_bold else "Lato Light",
                             size_pt=float(prim.get("size", 10)), bold=first_col_bold)
                set_run_scheme_color(run, accent_color if first_col_bold else "navy")
    return frame


def _render_primitives(shapes_host, primitives, accent, depth=0):
    for prim in primitives:
        ptype = prim.get("type")
        color = prim.get("fill") or prim.get("color") or accent
        if ptype == "arrow":
            shape_type = _ARROW_SHAPES.get(prim.get("direction", "right"), MSO_SHAPE.RIGHT_ARROW)
            shp = _shape(shapes_host, shape_type, prim["x"], prim["y"], prim["w"], prim["h"], fill=color,
                         line=prim.get("line"))
            if prim.get("rotation"):
                shp.rotation = float(prim["rotation"])
        elif ptype in _FREEFORM_SHAPES:
            shp = _shape(shapes_host, _FREEFORM_SHAPES[ptype], prim["x"], prim["y"], prim["w"], prim["h"],
                         fill=color, line=prim.get("line"))
            if prim.get("rotation"):
                shp.rotation = float(prim["rotation"])
            if prim.get("text"):
                tf = shp.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = str(prim["text"])
                set_run_font(run, prim.get("font", "Gotham Bold"),
                             size_pt=float(prim.get("size", 12)), bold=bool(prim.get("bold", True)))
                set_run_scheme_color(run, prim.get("text_color", "white"))
        elif ptype == "line":
            connector = shapes_host.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, _emu(prim["x1"]), _emu(prim["y1"]), _emu(prim["x2"]), _emu(prim["y2"]))
            set_scheme_line(connector, color, width_pt=float(prim.get("width_pt", 1.2)))
        elif ptype == "connector":
            _freeform_connector(shapes_host, prim, accent)
        elif ptype == "textbox":
            _freeform_textbox(shapes_host, prim, accent)
        elif ptype == "icon":
            icons.draw_icon(shapes_host, prim.get("icon"), prim["x"], prim["y"], prim["w"],
                             color=prim.get("color", accent), bg=prim.get("bg", "white"))
        elif ptype == "table":
            _freeform_table(shapes_host, prim, accent)
        elif ptype == "image":
            path = Path(prim.get("path", ""))
            if path.exists():
                h = prim.get("h")
                shapes_host.shapes.add_picture(str(path), _emu(prim["x"]), _emu(prim["y"]),
                                               width=_emu(prim["w"]), height=_emu(h) if h else None)
        elif ptype == "group" and depth < 2:
            group = shapes_host.shapes.add_group_shape()
            _render_primitives(group, prim.get("primitives", []), accent, depth + 1)


def _render_custom_freeform(slide, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    _render_primitives(slide, fill.get("primitives", []), accent)


_CHART_TYPE_MAP_NAMES = {
    "column": "COLUMN_CLUSTERED",
    "stacked-column": "COLUMN_STACKED",
    "bar": "BAR_CLUSTERED",
    "stacked-bar": "BAR_STACKED",
    "line": "LINE_MARKERS",
    "area": "AREA",
    "pie": "PIE",
    "donut": "DOUGHNUT",
    "scatter": "XY_SCATTER_LINES",
}
_DEFAULT_SERIES_ACCENTS = ["orange", "royal", "emerald", "grapefruit", "sky", "navy"]


def _render_chart_native(slide, comp, slide_spec, accent):
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    fill = slide_spec.get("content_schema_fill", {})
    visual = _visual(comp, slide_spec, {})
    chart_kind = fill.get("chart_type", "column")
    xl_type = getattr(XL_CHART_TYPE, _CHART_TYPE_MAP_NAMES[chart_kind])
    series_specs = fill.get("series", [])

    if chart_kind == "scatter":
        data = XyChartData()
        for s in series_specs:
            xy = data.add_series(s.get("name", ""))
            for xv, yv in zip(s.get("x_values", []), s.get("y_values", [])):
                xy.add_data_point(xv, yv)
    else:
        data = CategoryChartData()
        data.categories = [str(c) for c in fill.get("categories", [])]
        for s in series_specs:
            data.add_series(s.get("name", ""), tuple(float(v) for v in s.get("values", [])))

    x = float(visual.get("x", 0.95)); y = float(visual.get("y", 1.95))
    w = float(visual.get("width", 11.4)); h = float(visual.get("height", 4.6))
    frame = slide.shapes.add_chart(xl_type, _emu(x), _emu(y), _emu(w), _emu(h), data)
    chart = frame.chart
    chart.has_title = False

    show_legend = bool(visual.get("legend", len(series_specs) > 1 or chart_kind in ("pie", "donut")))
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Lato Light"
        chart.legend.font.size = Pt(11)

    if chart_kind in ("pie", "donut"):
        slice_colors = fill.get("slice_colors") or _DEFAULT_SERIES_ACCENTS
        plot_series = chart.plots[0].series[0]
        n_points = len(fill.get("categories", []))
        for i in range(n_points):
            set_fillformat_scheme(plot_series.points[i].format.fill, slice_colors[i % len(slice_colors)])
    else:
        for i, plot_series in enumerate(chart.series):
            color = _color(series_specs[i].get("color") if i < len(series_specs) else None,
                           None) or (accent if i == 0 else _DEFAULT_SERIES_ACCENTS[i % len(_DEFAULT_SERIES_ACCENTS)])
            if chart_kind in ("line", "scatter"):
                set_scheme_line(plot_series.format, color, width_pt=2.25)
            else:
                set_fillformat_scheme(plot_series.format.fill, color)

    if chart_kind not in ("pie", "donut", "scatter"):
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.name = "Lato Light"
        cat_axis.tick_labels.font.size = Pt(10.5)
        val_axis = chart.value_axis
        val_axis.tick_labels.font.name = "Lato Light"
        val_axis.tick_labels.font.size = Pt(10.5)
        val_axis.has_major_gridlines = True
        if visual.get("number_format"):
            val_axis.tick_labels.number_format = str(visual["number_format"])
            val_axis.tick_labels.number_format_is_linked = False

    if visual.get("data_labels") and chart_kind not in ("scatter",):
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.name = "Gotham Bold"
        plot.data_labels.font.size = Pt(10)
        if visual.get("number_format"):
            plot.data_labels.number_format = str(visual["number_format"])
            plot.data_labels.number_format_is_linked = False
    return frame


def _build_custom_slide(prs, comp, slide_spec, accent):
    fill = slide_spec.get("content_schema_fill", {})
    if comp["id"] in FULL_BLEED_COMPONENTS:
        slide = prs.slides.add_slide(find_layout(prs, "Leeg"))
    else:
        slide = prs.slides.add_slide(find_layout(prs, CONTENT_LAYOUT))
        _fill_content_frame(slide, fill)

    if comp["id"] == "kpi-trio":
        _render_kpi_trio(slide, comp, slide_spec, accent)

    elif comp["id"] == "content-cards":
        _render_content_cards(slide, comp, slide_spec, accent)

    elif comp["id"] == "chart-static":
        _render_chart_static(slide, comp, slide_spec, accent)
    elif comp["id"] == "process-timeline":
        _render_process_timeline(slide, comp, slide_spec, accent)
    elif comp["id"] == "schema-grid":
        _render_schema_grid(slide, comp, slide_spec, accent)
    elif comp["id"] == "image-icon-trio":
        _render_image_icon_trio(slide, comp, slide_spec, accent)
    elif comp["id"] == "matrix-2x2":
        _render_matrix_2x2(slide, comp, slide_spec, accent)
    elif comp["id"] == "layer-stack":
        _render_layer_stack(slide, comp, slide_spec, accent)
    elif comp["id"] == "cycle-diagram":
        _render_cycle_diagram(slide, comp, slide_spec, accent)
    elif comp["id"] == "scenario-cards":
        _render_scenario_cards(slide, comp, slide_spec, accent)
    elif comp["id"] == "assessment-table":
        _render_assessment_table(slide, comp, slide_spec, accent)
    elif comp["id"] == "mechanism-diagram":
        _render_mechanism_diagram(slide, comp, slide_spec, accent)
    elif comp["id"] == "divider-block":
        _render_divider_block(slide, comp, slide_spec, accent)
    elif comp["id"] == "stat-banner":
        _render_stat_banner(slide, comp, slide_spec, accent)
    elif comp["id"] == "swimlane-columns":
        _render_swimlane_columns(slide, comp, slide_spec, accent)
    elif comp["id"] == "closing-geometric":
        _render_closing_geometric(slide, comp, slide_spec, accent)
    elif comp["id"] == "custom-freeform":
        _render_custom_freeform(slide, comp, slide_spec, accent)
    elif comp["id"] == "chart-native":
        _render_chart_native(slide, comp, slide_spec, accent)
    return slide


def build_deck(spec: dict, out_path) -> Path:
    errors = validate_spec(spec)
    if errors:
        raise SpecError("invalid deck-spec:\n  - " + "\n  - ".join(errors))
    comps = load_components()
    meta = spec.get("meta") or {}
    deck_accent = meta.get("accent", "orange")
    prs = load_template_presentation()
    for slide_spec in spec["slides"]:
        comp = comps[slide_spec["component_id"]]
        accent = resolve_accent(meta, slide_spec.get("category"), deck_accent)
        if comp["renderer"] == "template":
            _build_template_slide(prs, comp, slide_spec)
        else:
            _build_custom_slide(prs, comp, slide_spec, accent)
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    import sys, json
    spec = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    dest = sys.argv[2] if len(sys.argv) > 2 else (spec.get("meta", {}).get("output") or "output/deck.pptx")
    print("built", build_deck(spec, dest))
