"""Cheap brand/text QA over a built deck (no rendering). v2 for the html2pptx
pipeline: slides are plain text boxes/shapes; brand hex is by design and the
allowed set comes from engine/web/tokens.json (palette + precomputed tints)."""
from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ENGINE = Path(__file__).resolve().parents[1]

ALLOWED_FONTS = {"Montserrat Light", "Lato Light", "Gotham Bold"}
LEFTOVER_MARKERS = (
    "click to edit", "tijdelijke aanduiding", "text placeholder", "lorem ipsum",
    "vervang deze", "action title in all caps", "optionele subtitel",
)
NEUTRAL_HEX = {"FFFFFF", "FEFFFF", "000000"}
TITLE_TOP_MAX = Inches(0.9)
TITLE_MIN_PT = 14


def _allowed_hexes() -> set[str]:
    tokens = json.loads((ENGINE / "web" / "tokens.json").read_text(encoding="utf-8"))
    return {h.upper() for h in tokens["allowed_hex"]} | NEUTRAL_HEX


def qa_deck(pptx_path) -> dict:
    prs = Presentation(str(pptx_path))
    allowed = _allowed_hexes()
    findings = []
    for si, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text
            texts.append(t)
            fonts, sizes = set(), []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
                    if run.font.size is not None:
                        sizes.append(run.font.size.pt)
            for fn in sorted(fonts - ALLOWED_FONTS):
                findings.append({"slide": si, "axis": "Design", "severity": "warn",
                                 "message": f"non-brand font {fn!r}"})
            is_title = ("Gotham Bold" in fonts and shape.top is not None
                        and shape.top < TITLE_TOP_MAX
                        and any(s >= TITLE_MIN_PT for s in sizes))
            if is_title and t.strip() and t != t.upper():
                findings.append({"slide": si, "axis": "Design", "severity": "critical",
                                 "message": f"titel niet in ALL CAPS: {t!r}"})
            low = t.lower()
            for marker in LEFTOVER_MARKERS:
                if marker in low:
                    findings.append({"slide": si, "axis": "Content", "severity": "critical",
                                     "message": f"leftover scaffold/placeholder text: {t!r}"})
        if not any(t.strip() for t in texts):
            findings.append({"slide": si, "axis": "Content", "severity": "warn",
                             "message": "slide has no text"})
        xml = slide._element.xml
        for hexv in set(re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml)):
            if hexv.upper() not in allowed:
                findings.append({"slide": si, "axis": "Design", "severity": "warn",
                                 "message": f"off-brand color #{hexv.upper()}"})
    critical = sum(1 for f in findings if f["severity"] == "critical")
    return {"findings": findings, "critical": critical}


if __name__ == "__main__":
    import sys
    print(json.dumps(qa_deck(sys.argv[1]), indent=2, ensure_ascii=False))
