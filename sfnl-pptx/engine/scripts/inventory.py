"""Extract a structural + text inventory from an existing .pptx for sfnl-deck-edit.

Written for editing a .pptx that has no HTML/deck.json source (see
skills/sfnl-deck-edit/SKILL.md). Reports, per slide, every text-bearing shape's position,
paragraph text/formatting, whether it hangs off the slide edge, and whether it overlaps
another shape on the same slide. Does NOT estimate in-shape text overflow from font
metrics — that needs real text layout to be trustworthy, and this plugin already has that
via scripts.render (actual PowerPoint rendering) plus the mandatory visual-inspection step.

Usage:
    python -m scripts.inventory input.pptx inventory.json
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.util import Emu

PT_PER_EMU = 1 / 12700
SLIDE_NUMBER_PLACEHOLDER = "SLIDE_NUMBER"


def _pt(value: Emu | int) -> float:
    return round(int(value) * PT_PER_EMU, 2)


def _placeholder_type(shape: BaseShape) -> str | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    fmt = shape.placeholder_format
    if fmt is None or fmt.type is None:
        return None
    return str(fmt.type).split(".")[-1].split(" ")[0]


def _paragraph_bullet(paragraph) -> tuple[bool, int | None]:
    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        return False, None
    has_bullet = pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None
    if not has_bullet:
        return False, None
    return True, paragraph.level or 0


def _paragraph_alignment(paragraph) -> str | None:
    if paragraph.alignment is None:
        return None
    name = str(paragraph.alignment).split(" ")[0].split(".")[-1]
    return None if name == "LEFT" else name


def _run_font(paragraph) -> dict[str, Any]:
    if not paragraph.runs:
        return {}
    font = paragraph.runs[0].font
    data: dict[str, Any] = {}
    if font.name:
        data["font_name"] = font.name
    if font.size is not None:
        data["font_size_pt"] = round(font.size.pt, 1)
    if font.bold is not None:
        data["bold"] = font.bold
    if font.italic is not None:
        data["italic"] = font.italic
    if font.underline is not None:
        data["underline"] = bool(font.underline)
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            data["color"] = str(font.color.rgb)
    except (AttributeError, TypeError):
        pass
    if "color" not in data:
        try:
            if font.color and font.color.theme_color is not None:
                data["theme_color"] = str(font.color.theme_color).split(".")[-1].split(" ")[0]
        except (AttributeError, TypeError):
            pass
    return data


def _paragraphs(shape: BaseShape) -> list[dict[str, Any]]:
    out = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        entry: dict[str, Any] = {"text": text}
        bullet, level = _paragraph_bullet(paragraph)
        if bullet:
            entry["bullet"] = True
            entry["level"] = level
        alignment = _paragraph_alignment(paragraph)
        if alignment:
            entry["alignment"] = alignment
        entry.update(_run_font(paragraph))
        out.append(entry)
    return out


def _is_content_shape(shape: BaseShape) -> bool:
    if not shape.has_text_frame:
        return False
    if not shape.text_frame.text.strip():
        return False
    ptype = _placeholder_type(shape)
    if ptype == SLIDE_NUMBER_PLACEHOLDER:
        return False
    if ptype == "FOOTER" and shape.text_frame.text.strip().isdigit():
        return False
    return True


def _group_child_space(shape) -> tuple[int, int, int, int] | None:
    """Return (chOff_x, chOff_y, chExt_cx, chExt_cy) for a GroupShape, or None if absent.

    A group's children are positioned in a local "child coordinate space" (chOff/chExt)
    that maps onto the group's own on-slide box (off/ext) — not necessarily 1:1, so a
    child's raw .left/.top can't just be added to the group's .left/.top.
    """
    xfrm = shape._element.find(f"{qn('p:grpSpPr')}/{qn('a:xfrm')}")
    if xfrm is None:
        return None
    ch_off = xfrm.find(qn("a:chOff"))
    ch_ext = xfrm.find(qn("a:chExt"))
    if ch_off is None or ch_ext is None:
        return None
    return (
        int(ch_off.get("x")), int(ch_off.get("y")),
        int(ch_ext.get("cx")), int(ch_ext.get("cy")),
    )


def _flatten_shapes(shapes, tx: float = 0, ty: float = 0, sx: float = 1, sy: float = 1):
    """Yield (shape, absolute_left_emu, absolute_top_emu, absolute_width_emu,
    absolute_height_emu), recursing into groups.

    (tx, ty, sx, sy) is the affine transform mapping this shape-list's local coordinate
    space onto the slide: absolute = (tx + sx * local_x, ty + sy * local_y). Composed
    across nested groups via each group's own off/ext -> chOff/chExt mapping.
    """
    for shape in shapes:
        abs_left = tx + sx * (shape.left or 0)
        abs_top = ty + sy * (shape.top or 0)
        abs_w = sx * (shape.width or 0)
        abs_h = sy * (shape.height or 0)
        if hasattr(shape, "shapes"):  # GroupShape — recurse with a composed child transform
            child_space = _group_child_space(shape)
            if child_space is None:
                yield from _flatten_shapes(shape.shapes, abs_left, abs_top, sx, sy)
                continue
            ch_x, ch_y, ch_w, ch_h = child_space
            child_sx = abs_w / ch_w if ch_w else sx
            child_sy = abs_h / ch_h if ch_h else sy
            child_tx = abs_left - child_sx * ch_x
            child_ty = abs_top - child_sy * ch_y
            yield from _flatten_shapes(shape.shapes, child_tx, child_ty, child_sx, child_sy)
            continue
        yield shape, abs_left, abs_top, abs_w, abs_h


def _rect_overlap_area_pt2(a: dict, b: dict) -> float:
    left = max(a["left_pt"], b["left_pt"])
    top = max(a["top_pt"], b["top_pt"])
    right = min(a["left_pt"] + a["width_pt"], b["left_pt"] + b["width_pt"])
    bottom = min(a["top_pt"] + a["height_pt"], b["top_pt"] + b["height_pt"])
    if right <= left or bottom <= top:
        return 0.0
    return round((right - left) * (bottom - top), 1)


def iter_content_shapes(prs: Presentation):
    """Yield (slide_index, shape_id, shape, abs_left_emu, abs_top_emu, abs_width_emu,
    abs_height_emu) for every text-bearing content shape. The abs_* values account for
    group-shape nesting (a shape's own .left/.top/.width/.height are relative to its
    immediate parent group's coordinate space, not the slide).

    shape_id is stable for a given slide's shape order (flattened groups, document order) —
    shared by inventory.py and replace.py so both agree on which shape is which.
    """
    for slide_index, slide in enumerate(prs.slides):
        shape_id = 0
        for shape, abs_left, abs_top, abs_w, abs_h in _flatten_shapes(slide.shapes):
            if not _is_content_shape(shape):
                continue
            yield slide_index, shape_id, shape, abs_left, abs_top, abs_w, abs_h
            shape_id += 1


def build_inventory(pptx_path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_w, slide_h = _pt(prs.slide_width), _pt(prs.slide_height)
    by_slide: dict[int, list[dict[str, Any]]] = {}

    for slide_index, shape_id, shape, abs_left, abs_top, abs_w, abs_h in iter_content_shapes(prs):
        rect = {
            "left_pt": _pt(abs_left),
            "top_pt": _pt(abs_top),
            "width_pt": _pt(abs_w),
            "height_pt": _pt(abs_h),
        }
        entry: dict[str, Any] = {"id": shape_id, "name": shape.name, **rect}
        ptype = _placeholder_type(shape)
        if ptype:
            entry["placeholder_type"] = ptype

        off_right = rect["left_pt"] + rect["width_pt"] - slide_w
        off_bottom = rect["top_pt"] + rect["height_pt"] - slide_h
        off_slide = {}
        if off_right > 0.5:
            off_slide["right_pt"] = round(off_right, 1)
        if off_bottom > 0.5:
            off_slide["bottom_pt"] = round(off_bottom, 1)
        if off_slide:
            entry["off_slide"] = off_slide

        entry["paragraphs"] = _paragraphs(shape)
        by_slide.setdefault(slide_index, []).append(entry)

    slides_out = []
    for slide_index, shapes_out in by_slide.items():
        for i, a in enumerate(shapes_out):
            for b in shapes_out[i + 1 :]:
                area = _rect_overlap_area_pt2(a, b)
                if area > 0:
                    a.setdefault("overlaps", []).append({"shape_id": b["id"], "area_pt2": area})
                    b.setdefault("overlaps", []).append({"shape_id": a["id"], "area_pt2": area})
        slides_out.append({"index": slide_index, "shapes": shapes_out})

    return {
        "source": str(pptx_path),
        "slide_width_pt": slide_w,
        "slide_height_pt": slide_h,
        "slide_count": len(prs.slides),
        "slides": slides_out,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("output", help="Output inventory .json file")
    args = parser.parse_args()

    inventory = build_inventory(args.input)
    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(inventory, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Wrote {out_path} — {len(inventory['slides'])} slide(s) with content")


if __name__ == "__main__":
    main()
