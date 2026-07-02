import json
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from scripts.build_from_spec import build_deck

FIX = Path(__file__).resolve().parent / "fixtures" / "sample_visual_spec.json"


def _build(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "visual.pptx")
    from pptx import Presentation

    return Presentation(str(out))


def _auto_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _line_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]


def _slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)


def _near(actual, expected, tolerance=Inches(0.06)):
    return abs(int(actual) - int(expected)) <= int(tolerance)


def test_content_cards_render_panels_icons_and_respect_geometry(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[0]

    assert len(_auto_shapes(slide)) >= 6
    assert "✓" in _slide_text(slide)
    assert "⚙" in _slide_text(slide)
    assert "♡" in _slide_text(slide)
    assert any(
        _near(shape.left, Inches(1.1))
        and _near(shape.top, Inches(2.0))
        and _near(shape.width, Inches(2.8))
        for shape in _auto_shapes(slide)
    )


def test_kpi_trio_renders_panels_and_progress_bar(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[1]

    assert len(_auto_shapes(slide)) >= 8
    assert "36% van aflossingen terugontvangen" in _slide_text(slide)
    assert any(_near(shape.left, Inches(0.65)) for shape in _auto_shapes(slide))
    assert "srgbClr" not in slide._element.xml


def test_process_timeline_renders_arrow_shapes_with_step_icons(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[2]

    assert len(_auto_shapes(slide)) >= 4
    for number in ["1", "2", "3", "4"]:
        assert number in _slide_text(slide)
    assert any(_near(shape.left, Inches(0.85)) for shape in _auto_shapes(slide))


def test_schema_grid_renders_boxes_and_connectors(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[3]

    assert len(_auto_shapes(slide)) >= 8
    assert len(_line_shapes(slide)) >= 3
    assert "Team HR" in _slide_text(slide)
    assert "Communicatie" in _slide_text(slide)
    assert "srgbClr" not in slide._element.xml

