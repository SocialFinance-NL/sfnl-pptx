# sfnl-pptx/tests/test_build.py
from pathlib import Path
from pptx import Presentation
from scripts.build_from_spec import build_deck
import json

FIX = Path(__file__).resolve().parent / "fixtures" / "sample_spec.json"

def test_build_produces_openable_deck(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "demo.pptx")
    assert out.exists()
    prs = Presentation(str(out))           # re-opens without error => not corrupt
    assert len(prs.slides) == 3

def test_title_text_lands_in_placeholder(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "demo.pptx")
    prs = Presentation(str(out))
    text = "\n".join(sh.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "MAATSCHAPPELIJKE WAARDE, MEETBAAR GEMAAKT" in text  # titels altijd ALL CAPS

def test_custom_slide_uses_template_content_layout(tmp_path):
    """Custom slides bouwen op de sjabloonlayout (wit vlak, oranje streepje, titelplaceholder)."""
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "demo.pptx")
    prs = Presentation(str(out))
    kpi_slide = prs.slides[1]
    assert kpi_slide.slide_layout.name == "Titel, subtitel"
    titles = [ph for ph in kpi_slide.placeholders if ph.placeholder_format.idx == 0]
    assert titles and titles[0].text_frame.text == "IMPACT IN CIJFERS"

def test_custom_kpi_slide_has_no_hardcoded_hex(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "demo.pptx")
    prs = Presentation(str(out))
    kpi_xml = prs.slides[1]._element.xml
    assert "3,2x" in kpi_xml
    # custom shapes must color via schemeClr, never srgbClr
    assert "srgbClr" not in kpi_xml
