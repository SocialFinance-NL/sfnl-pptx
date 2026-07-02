import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches

from scripts.build_from_spec import build_deck

FIX = Path(__file__).resolve().parent / "fixtures" / "sample_reference_style_spec.json"


def _build(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "reference-style.pptx")
    return Presentation(str(out))


def _auto_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _line_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]


def _slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)


def _near(actual, expected, tolerance=Inches(0.06)):
    return abs(int(actual) - int(expected)) <= int(tolerance)


def test_divider_block_is_full_bleed_and_uses_mapped_category_accent(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[0]
    panels = [s for s in _auto_shapes(slide) if _near(s.width, Inches(5.3)) and _near(s.height, Inches(7.5))]
    assert panels, "expected a full-height color panel sized to panel_width"
    assert 'val="accent2"' in panels[0].fill._xPr.xml  # category "vraagstuk" -> grapefruit
    assert "HET VRAAGSTUK" in _slide_text(slide)  # titels altijd ALL CAPS


def test_stat_banner_renders_dark_bars_with_values(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[1]
    bars = [s for s in _auto_shapes(slide) if _near(s.height, Inches(1.15))]
    assert len(bars) >= 2
    text = _slide_text(slide)
    assert "593" in text
    assert "1.240" in text


def test_swimlane_columns_render_color_coded_headers_and_items(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[2]
    text = _slide_text(slide)
    for heading in ["Vraagstuk", "Activiteiten", "Outcomes", "Impact"]:
        assert heading in text
    assert "Eenzaamheid" in text
    header_colors = set()
    for shape in _auto_shapes(slide):
        xml = shape.fill._xPr.xml
        for slot in ("accent1", "accent2", "accent3", "accent4", "accent5"):
            if f'val="{slot}"' in xml:
                header_colors.add(slot)
    assert len(header_colors) >= 3, "swimlane headers should use distinct per-column colors"


def test_content_cards_pastel_tint_variant_uses_tint_not_hex(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[3]
    xml = slide._element.xml
    assert "lumOff" in xml, "pastel-tint variant should tint card backgrounds via lumMod/lumOff"
    assert "srgbClr" not in xml


def test_process_timeline_chevron_variant_uses_chevron_shape(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[4]
    chevrons = [s for s in _auto_shapes(slide) if getattr(s, "auto_shape_type", None) == MSO_SHAPE.CHEVRON]
    assert len(chevrons) >= 2


def test_custom_freeform_renders_declared_primitives(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[5]
    assert len(_line_shapes(slide)) >= 1
    assert "Vrije tekst" in _slide_text(slide)
    assert len(_auto_shapes(slide)) >= 4  # rect + oval + icon shapes (gear draws several)


def test_closing_geometric_is_full_bleed_with_geometric_shapes(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[6]
    panels = [s for s in _auto_shapes(slide) if _near(s.width, Inches(13.33)) and _near(s.height, Inches(7.5))]
    assert panels, "expected a full-bleed background panel"
    triangles = [s for s in _auto_shapes(slide) if getattr(s, "auto_shape_type", None) == MSO_SHAPE.ISOSCELES_TRIANGLE]
    assert len(triangles) >= 2
    text = _slide_text(slide)
    assert "DANK U WEL" in text  # titels altijd ALL CAPS
