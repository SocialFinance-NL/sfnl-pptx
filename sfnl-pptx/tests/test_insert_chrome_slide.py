import pytest
from pptx import Presentation
from pptx.util import Pt

from scripts.insert_chrome_slide import insert_chrome_slide


def _deck(tmp_path):
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(405)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Pt(20), Pt(20), Pt(200), Pt(50)).text_frame.text = "BESTAAND"
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_insert_chrome_slide_adds_full_bleed_art_and_slot_text(tmp_path):
    out = tmp_path / "out.pptx"
    insert_chrome_slide(
        _deck(tmp_path),
        out,
        "divider-01",
        0,
        {"title": "NIEUWE SECTIE", "body": "Introductietekst"},
    )

    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    first = prs.slides[0]
    texts = [shape.text_frame.text for shape in first.shapes if shape.has_text_frame]
    assert texts == ["NIEUWE SECTIE", "Introductietekst"]
    picture = first.shapes[0]
    assert picture.left == 0
    assert picture.top == 0
    assert picture.width == prs.slide_width
    assert picture.height == prs.slide_height
    assert prs.slides[1].shapes[0].text_frame.text == "BESTAAND"


def test_insert_chrome_slide_supports_keys_without_text_slots(tmp_path):
    out = tmp_path / "out.pptx"
    insert_chrome_slide(_deck(tmp_path), out, "cover-02", 1)
    assert len(Presentation(str(out)).slides[1].shapes) == 1


@pytest.mark.parametrize(
    ("key", "position", "texts", "match"),
    [
        ("divider-99", 0, {}, "Unknown chrome key"),
        ("divider-01", 99, {}, "out of range"),
        ("divider-01", 0, {"subtitle": "niet geldig"}, "has no slot"),
    ],
)
def test_insert_chrome_slide_fails_loudly_on_invalid_requests(tmp_path, key, position, texts, match):
    with pytest.raises(ValueError, match=match):
        insert_chrome_slide(_deck(tmp_path), tmp_path / "out.pptx", key, position, texts)
