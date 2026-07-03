from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from scripts.inventory import build_inventory


def _save_deck(tmp_path, build):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build(slide)
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_inventory_extracts_text_geometry_formatting_and_issues(tmp_path):
    def build(slide):
        first = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.7))
        run = first.text_frame.paragraphs[0].add_run()
        run.text = "Eerste tekst"
        run.font.name = "Gotham Bold"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("201B5C")

        overlap = slide.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(3), Inches(0.7))
        overlap.text_frame.text = "Overlap"

        off_slide = slide.shapes.add_textbox(Inches(9.5), Inches(5.2), Inches(1), Inches(0.8))
        off_slide.text_frame.text = "Buiten canvas"

    inventory = build_inventory(_save_deck(tmp_path, build))

    assert inventory["slide_count"] == 1
    shapes = inventory["slides"][0]["shapes"]
    assert [shape["id"] for shape in shapes] == [0, 1, 2]
    assert shapes[0]["paragraphs"][0] | {"text": "Eerste tekst"} == shapes[0]["paragraphs"][0]
    assert shapes[0]["paragraphs"][0]["font_name"] == "Gotham Bold"
    assert shapes[0]["paragraphs"][0]["font_size_pt"] == 18
    assert shapes[0]["paragraphs"][0]["bold"] is True
    assert shapes[0]["paragraphs"][0]["color"] == "201B5C"
    assert shapes[0]["overlaps"][0]["shape_id"] == 1
    assert shapes[2]["off_slide"]["right_pt"] > 0
    assert shapes[2]["off_slide"]["bottom_pt"] > 0
