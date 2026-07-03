import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from scripts.inventory import build_inventory
from scripts.replace import apply_replacements


def _deck(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for idx, text in enumerate(("Oud 1", "Oud 2")):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5 + idx), Inches(4), Inches(0.6))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.name = "Lato Light"
        run.font.size = Pt(10)
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_replace_applies_plan_and_clears_unmentioned_content_shapes(tmp_path):
    source = _deck(tmp_path)
    out = tmp_path / "edited.pptx"
    apply_replacements(
        source,
        {
            "slides": [
                {
                    "index": 0,
                    "shapes": [
                        {
                            "id": 0,
                            "paragraphs": [
                                {
                                    "text": "Nieuwe tekst",
                                    "font_name": "Gotham Bold",
                                    "font_size_pt": 18,
                                    "bold": True,
                                    "color": "F87F4F",
                                }
                            ],
                        }
                    ],
                }
            ]
        },
        out,
    )

    inventory = build_inventory(out)
    shapes = inventory["slides"][0]["shapes"]
    assert len(shapes) == 1
    paragraph = shapes[0]["paragraphs"][0]
    assert paragraph["text"] == "Nieuwe tekst"
    assert paragraph["font_name"] == "Gotham Bold"
    assert paragraph["font_size_pt"] == 18
    assert paragraph["bold"] is True
    assert paragraph["color"] == "F87F4F"

    prs = Presentation(str(out))
    second_shape = prs.slides[0].shapes[1]
    assert second_shape.text_frame.text == ""


def test_replace_fails_on_unknown_inventory_shape(tmp_path):
    with pytest.raises(ValueError, match="not found by inventory"):
        apply_replacements(
            _deck(tmp_path),
            {"slides": [{"index": 0, "shapes": [{"id": 99, "paragraphs": [{"text": "x"}]}]}]},
            tmp_path / "edited.pptx",
        )
