import json
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from scripts.build_from_spec import build_deck

FIX = Path(__file__).resolve().parent / "fixtures" / "sample_html_components_spec.json"


def _build(tmp_path):
    spec = json.loads(FIX.read_text(encoding="utf-8"))
    out = build_deck(spec, tmp_path / "html-components.pptx")
    from pptx import Presentation

    return Presentation(str(out))


def _auto_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _line_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]


def _text(slide):
    return "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)


def _near(actual, expected, tolerance=Inches(0.08)):
    return abs(int(actual) - int(expected)) <= int(tolerance)


def test_matrix_2x2_renders_quadrants_axes_and_respects_geometry(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[0]

    assert len(_auto_shapes(slide)) >= 8
    assert "Strategisch" in _text(slide)
    assert "Vermijden" in _text(slide)
    assert "Effect" in _text(slide)
    assert "Risico" in _text(slide)
    assert any(_near(shape.left, Inches(0.95)) and _near(shape.top, Inches(1.95)) for shape in _auto_shapes(slide))


def test_layer_stack_renders_indented_layers(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[1]
    shapes = _auto_shapes(slide)
    layer_shapes = [shape for shape in shapes if _near(shape.height, Inches(0.55))]

    assert len(layer_shapes) >= 5
    assert "Niveau 5" in _text(slide)
    assert "Operationele registratie" in _text(slide)
    assert min(shape.left for shape in layer_shapes) < max(shape.left for shape in layer_shapes)


def test_cycle_diagram_renders_center_nodes_and_connectors(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[2]

    assert len(_auto_shapes(slide)) >= 5
    assert len(_line_shapes(slide)) >= 4
    assert "Continu" in _text(slide)
    assert "Bijsturen" in _text(slide)


def test_scenario_cards_render_sections_and_bullets(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[3]

    assert len(_auto_shapes(slide)) >= 9
    assert "Scenario A" in _text(slide)
    assert "Regionaal model" in _text(slide)
    assert "Governance-uitdaging" in _text(slide)


def test_assessment_table_renders_rows_scores_and_headers(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[4]

    assert len(_auto_shapes(slide)) >= 8
    assert "Criterium" in _text(slide)
    assert "Toegevoegde waarde" in _text(slide)
    assert "Middel" in _text(slide)
    assert "Hoog" in _text(slide)


def test_mechanism_diagram_renders_nodes_flows_and_labels(tmp_path):
    prs = _build(tmp_path)
    slide = prs.slides[5]

    assert len(_auto_shapes(slide)) >= 10
    assert len(_line_shapes(slide)) >= 4
    assert "Maatwerkbudget" in _text(slide)
    assert "Learning partner" in _text(slide)
    assert "Evidence" in _text(slide)
    assert "Budget" in _text(slide)

