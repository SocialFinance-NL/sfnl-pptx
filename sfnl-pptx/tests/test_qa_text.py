"""qa_text v2: rules over html2pptx-built decks (plain text boxes, brand hex by design)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from scripts.qa_text import qa_deck


def _deck(tmp_path, build):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    build(slide)
    p = tmp_path / "deck.pptx"
    prs.save(p)
    return p


def _textbox(slide, text, font, size, top=Inches(0.4), color=None):
    box = slide.shapes.add_textbox(Inches(0.4), top, Inches(6), Inches(0.5))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    return box


def test_lowercase_gotham_title_is_critical(tmp_path):
    p = _deck(tmp_path, lambda s: _textbox(s, "geen caps titel", "Gotham Bold", 18))
    r = qa_deck(p)
    assert any("ALL CAPS" in f["message"] for f in r["findings"])
    assert r["critical"] >= 1


def test_caps_title_and_brand_colors_pass(tmp_path):
    def build(s):
        _textbox(s, "NETTE TITEL", "Gotham Bold", 18, color="201B5C")
        _textbox(s, "Body in Lato.", "Lato Light", 10, top=Inches(2))
    r = qa_deck(_deck(tmp_path, build))
    assert r["critical"] == 0


def test_big_number_lower_on_canvas_is_not_a_title(tmp_path):
    p = _deck(tmp_path, lambda s: _textbox(s, "31% minder", "Gotham Bold", 30, top=Inches(2.5)))
    assert qa_deck(p)["critical"] == 0


def test_non_brand_font_warns(tmp_path):
    p = _deck(tmp_path, lambda s: _textbox(s, "TEKST", "Comic Sans MS", 12))
    assert any("non-brand font" in f["message"] for f in qa_deck(p)["findings"])


def test_off_token_color_warns_but_tint_passes(tmp_path):
    import json
    tokens = json.loads((Path(__file__).resolve().parents[1] / "engine" / "web" / "tokens.json").read_text())
    a_tint = next(h for h in tokens["allowed_hex"] if h not in {"201B5C", "233348", "F87F4F"})
    def build(s):
        _textbox(s, "PAARS", "Lato Light", 10, color="8B00FF")
        _textbox(s, "TINT", "Lato Light", 10, top=Inches(2), color=a_tint)
    findings = qa_deck(_deck(tmp_path, build))["findings"]
    assert any("8B00FF" in f["message"] for f in findings)
    assert not any(a_tint in f["message"] for f in findings)


def test_leftover_scaffold_text_is_critical(tmp_path):
    p = _deck(tmp_path, lambda s: _textbox(s, "Vervang deze inhoud.", "Lato Light", 10, top=Inches(2)))
    r = qa_deck(p)
    assert any("leftover" in f["message"] for f in r["findings"])
    assert r["critical"] >= 1
