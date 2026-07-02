from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from scripts.colors import (
    set_scheme_fill, set_run_scheme_color, set_run_font, ACCENT_TO_SLOT,
    set_scheme_fill_tint, set_scheme_fill_shade, resolve_accent,
)

def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])

def test_accent_map_covers_brand_accents():
    assert ACCENT_TO_SLOT["orange"] == "accent1"
    assert ACCENT_TO_SLOT["emerald"] == "accent5"
    assert ACCENT_TO_SLOT["navy"] == "dk2"

def test_scheme_fill_uses_schemeclr_not_hex():
    slide = _blank_slide()
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    set_scheme_fill(shp, "emerald")
    xml = shp.fill._xPr.xml
    assert "schemeClr" in xml and 'val="accent5"' in xml
    assert "srgbClr" not in xml

def test_run_font_applies_name_and_size():
    slide = _blank_slide()
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Hi"
    set_run_font(run, "Gotham Bold", size_pt=28, bold=True)
    assert run.font.name == "Gotham Bold"
    assert run.font.size == Pt(28)
    assert run.font.bold is True


def test_scheme_fill_tint_uses_lummod_and_lumoff_no_hex():
    slide = _blank_slide()
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    set_scheme_fill_tint(shp, "emerald", 80)
    xml = shp.fill._xPr.xml
    assert 'val="accent5"' in xml
    assert 'lumMod val="20000"' in xml
    assert 'lumOff val="80000"' in xml
    assert "srgbClr" not in xml


def test_scheme_fill_shade_uses_lummod_only_no_hex():
    slide = _blank_slide()
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    set_scheme_fill_shade(shp, "navy", 25)
    xml = shp.fill._xPr.xml
    assert 'val="dk2"' in xml
    assert 'lumMod val="75000"' in xml
    assert "lumOff" not in xml
    assert "srgbClr" not in xml


def test_resolve_accent_falls_back_without_accent_map():
    assert resolve_accent({}, "vraagstuk", "orange") == "orange"
    assert resolve_accent({"accent": "orange"}, None, "orange") == "orange"


def test_resolve_accent_uses_category_mapping_when_present():
    meta = {"accent": "orange", "accent_map": {"vraagstuk": "grapefruit", "impact": "orange"}}
    assert resolve_accent(meta, "vraagstuk", "orange") == "grapefruit"
    assert resolve_accent(meta, "impact", "orange") == "orange"


def test_resolve_accent_falls_back_on_unmapped_category():
    meta = {"accent": "orange", "accent_map": {"vraagstuk": "grapefruit"}}
    assert resolve_accent(meta, "onbekend", "orange") == "orange"
