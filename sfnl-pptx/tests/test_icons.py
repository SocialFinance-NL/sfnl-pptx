from pptx import Presentation

from scripts.icons import ICON_NAMES, draw_icon


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_every_named_icon_draws_at_least_one_shape():
    for name in ICON_NAMES:
        slide = _blank_slide()
        before = len(slide.shapes)
        draw_icon(slide, name, 0.5, 0.5, 1.2, color="orange", bg="white")
        assert len(slide.shapes) > before, f"icon {name!r} drew no shapes"


def test_icons_use_schemeclr_never_hardcoded_hex():
    slide = _blank_slide()
    draw_icon(slide, "gear", 0.5, 0.5, 1.2, color="emerald", bg="white")
    assert "srgbClr" not in slide._element.xml
    assert "schemeClr" in slide._element.xml


def test_unknown_icon_name_falls_back_without_raising():
    slide = _blank_slide()
    before = len(slide.shapes)
    draw_icon(slide, "not-a-real-icon", 0.5, 0.5, 1.2, color="orange", bg="white")
    assert len(slide.shapes) > before
