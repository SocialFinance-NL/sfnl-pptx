import zipfile

from pptx import Presentation
from pptx.util import Inches

from ooxml.scripts.pack import pack
from ooxml.scripts.unpack import unpack
from ooxml.scripts.validate import validate


def _deck(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.6)).text_frame.text = (
        'Tekst met  spaties & "quotes"'
    )
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_ooxml_unpack_validate_pack_round_trip_preserves_openable_deck(tmp_path):
    source = _deck(tmp_path)
    workspace = tmp_path / "workspace"
    out = tmp_path / "roundtrip.pptx"

    unpack(source, workspace)
    assert validate(workspace) == []
    pack(workspace, out)

    reopened = Presentation(str(out))
    assert reopened.slides[0].shapes[0].text_frame.text == 'Tekst met  spaties & "quotes"'


def test_ooxml_validate_catches_broken_rels_reference(tmp_path):
    workspace = tmp_path / "workspace"
    unpack(_deck(tmp_path), workspace)
    rels_file = workspace / "ppt" / "slides" / "_rels" / "slide1.xml.rels"
    rels_file.write_text(
        rels_file.read_text(encoding="utf-8").replace("slideLayout", "missingSlideLayout"),
        encoding="utf-8",
    )

    assert any("broken reference" in error for error in validate(workspace))


def test_ooxml_validate_catches_malformed_xml(tmp_path):
    workspace = tmp_path / "workspace"
    unpack(_deck(tmp_path), workspace)
    slide_xml = workspace / "ppt" / "slides" / "slide1.xml"
    slide_xml.write_text(slide_xml.read_text(encoding="utf-8").replace("</p:sp>", "</p:sp", 1), encoding="utf-8")

    assert any("not well-formed" in error for error in validate(workspace))


def test_ooxml_pack_does_not_include_workspace_directory_name(tmp_path):
    workspace = tmp_path / "workspace"
    out = tmp_path / "roundtrip.pptx"
    unpack(_deck(tmp_path), workspace)
    pack(workspace, out)

    with zipfile.ZipFile(out) as archive:
        assert "[Content_Types].xml" in archive.namelist()
        assert not any(name.startswith("workspace/") for name in archive.namelist())
