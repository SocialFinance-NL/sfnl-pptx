"""Tests for the extended freeform primitive set and the chart-native component."""
import pytest

from scripts.spec import validate_spec
from scripts.build_from_spec import build_deck


def _spec(slides):
    return {"meta": {"title": "t", "lang": "nl", "accent": "emerald"}, "slides": slides}


def _freeform_slide(primitives, sid="s1"):
    return {"id": sid, "component_id": "custom-freeform", "action_title": "Bespoke exhibit draagt de boodschap",
            "content_schema_fill": {"title": "Bespoke", "primitives": primitives}}


def _chart_slide(fill, sid="s1"):
    base = {"title": "Chart"}
    base.update(fill)
    return {"id": sid, "component_id": "chart-native", "action_title": "Data draagt de boodschap",
            "content_schema_fill": base}


# ---------- validation: freeform ----------

def test_new_shape_types_validate():
    prims = [{"type": t, "x": 1.0, "y": 1.0, "w": 1.0, "h": 1.0, "fill": "sky"}
             for t in ("triangle", "diamond", "pentagon", "hexagon", "chevron", "donut", "pie")]
    assert validate_spec(_spec([_freeform_slide(prims)])) == []


def test_arrow_direction_validated():
    ok = [{"type": "arrow", "x": 1, "y": 1, "w": 1, "h": 1, "direction": "down"}]
    assert validate_spec(_spec([_freeform_slide(ok)])) == []
    bad = [{"type": "arrow", "x": 1, "y": 1, "w": 1, "h": 1, "direction": "diagonal"}]
    assert any("direction" in e for e in validate_spec(_spec([_freeform_slide(bad)])))


def test_connector_style_and_arrow_validated():
    ok = [{"type": "connector", "x1": 1, "y1": 1, "x2": 2, "y2": 2, "style": "elbow", "arrow": "both"}]
    assert validate_spec(_spec([_freeform_slide(ok)])) == []
    bad = [{"type": "connector", "x1": 1, "y1": 1, "x2": 2, "y2": 2, "style": "curvy"}]
    assert any("style" in e for e in validate_spec(_spec([_freeform_slide(bad)])))


def test_table_rows_validated():
    ok = [{"type": "table", "x": 1, "y": 1, "w": 5, "rows": [["a", "b"], ["c", "d"]]}]
    assert validate_spec(_spec([_freeform_slide(ok)])) == []
    ragged = [{"type": "table", "x": 1, "y": 1, "w": 5, "rows": [["a", "b"], ["c"]]}]
    assert any("same number" in e for e in validate_spec(_spec([_freeform_slide(ragged)])))
    empty = [{"type": "table", "x": 1, "y": 1, "w": 5, "rows": []}]
    assert any("rows" in e for e in validate_spec(_spec([_freeform_slide(empty)])))


def test_textbox_bullets_accepted_and_text_still_required():
    ok = [{"type": "textbox", "x": 1, "y": 1, "w": 3, "h": 1, "bullets": ["een", "twee"]}]
    assert validate_spec(_spec([_freeform_slide(ok)])) == []
    bad = [{"type": "textbox", "x": 1, "y": 1, "w": 3, "h": 1}]
    assert any("text" in e for e in validate_spec(_spec([_freeform_slide(bad)])))


def test_group_validates_children_and_depth():
    ok = [{"type": "group", "x": 1, "y": 1, "w": 3, "primitives": [
        {"type": "oval", "x": 1, "y": 1, "w": 1, "h": 1}]}]
    assert validate_spec(_spec([_freeform_slide(ok)])) == []
    bad_child = [{"type": "group", "x": 1, "y": 1, "w": 3, "primitives": [
        {"type": "blob", "x": 1, "y": 1, "w": 1, "h": 1}]}]
    assert any("blob" in e for e in validate_spec(_spec([_freeform_slide(bad_child)])))
    too_deep = [{"type": "group", "x": 1, "y": 1, "w": 3, "primitives": [
        {"type": "group", "x": 1, "y": 1, "w": 2, "primitives": [
            {"type": "group", "x": 1, "y": 1, "w": 1, "primitives": [
                {"type": "oval", "x": 1, "y": 1, "w": 1, "h": 1}]}]}]}]
    assert any("nest" in e for e in validate_spec(_spec([_freeform_slide(too_deep)])))


def test_image_requires_path():
    bad = [{"type": "image", "x": 1, "y": 1, "w": 3}]
    assert any("path" in e for e in validate_spec(_spec([_freeform_slide(bad)])))


def test_off_brand_color_still_rejected():
    bad = [{"type": "hexagon", "x": 1, "y": 1, "w": 1, "h": 1, "fill": "hotpink"}]
    assert any("hotpink" in e for e in validate_spec(_spec([_freeform_slide(bad)])))


# ---------- validation: chart-native ----------

def test_chart_native_valid_column():
    fill = {"chart_type": "column", "categories": ["a", "b"],
            "series": [{"name": "s", "values": [1, 2], "color": "emerald"}]}
    assert validate_spec(_spec([_chart_slide(fill)])) == []


def test_chart_native_rejects_unknown_type_and_bad_values():
    fill = {"chart_type": "radar", "categories": ["a"], "series": [{"name": "s", "values": [1]}]}
    assert any("chart_type" in e for e in validate_spec(_spec([_chart_slide(fill)])))
    fill = {"chart_type": "column", "categories": ["a"], "series": [{"name": "s", "values": ["x"]}]}
    assert any("numeric" in e for e in validate_spec(_spec([_chart_slide(fill)])))


def test_chart_native_requires_categories_and_series():
    fill = {"chart_type": "column", "series": [{"name": "s", "values": [1]}]}
    assert any("categories" in e for e in validate_spec(_spec([_chart_slide(fill)])))
    fill = {"chart_type": "column", "categories": ["a"], "series": []}
    assert any("series" in e for e in validate_spec(_spec([_chart_slide(fill)])))


def test_chart_native_scatter_needs_xy():
    fill = {"chart_type": "scatter", "series": [{"name": "s", "x_values": [1, 2], "y_values": [3, 4]}]}
    assert validate_spec(_spec([_chart_slide(fill)])) == []
    fill = {"chart_type": "scatter", "series": [{"name": "s", "values": [1, 2]}]}
    errors = validate_spec(_spec([_chart_slide(fill)]))
    assert any("x_values" in e for e in errors)


def test_chart_native_slice_colors_validated():
    fill = {"chart_type": "donut", "categories": ["a", "b"],
            "series": [{"name": "s", "values": [1, 2]}], "slice_colors": ["emerald", "neon"]}
    assert any("neon" in e for e in validate_spec(_spec([_chart_slide(fill)])))


# ---------- build ----------

def test_build_extended_freeform(tmp_path):
    prims = [
        {"type": "chevron", "x": 0.9, "y": 2.0, "w": 2.2, "h": 0.9, "fill": "emerald", "text": "Fase 1"},
        {"type": "arrow", "direction": "right", "x": 3.3, "y": 2.0, "w": 1.0, "h": 0.9, "fill": "sky"},
        {"type": "hexagon", "x": 4.6, "y": 1.85, "w": 1.3, "h": 1.2, "fill": "orange", "text": "Hub"},
        {"type": "connector", "style": "elbow", "arrow": "both", "x1": 6.0, "y1": 2.4, "x2": 8.0, "y2": 3.6,
         "color": "navy"},
        {"type": "donut", "x": 8.4, "y": 2.0, "w": 1.4, "h": 1.4, "fill": "royal"},
        {"type": "textbox", "x": 0.9, "y": 3.4, "w": 4.0, "h": 1.4, "bullets": ["Een", "Twee"], "size": 12},
        {"type": "table", "x": 5.4, "y": 4.4, "w": 6.5, "header": True, "first_col_bold": True,
         "rows": [["Criterium", "Score"], ["Draagvlak", "Hoog"]]},
        {"type": "group", "x": 0.9, "y": 5.2, "w": 3.0, "primitives": [
            {"type": "oval", "x": 0.9, "y": 5.2, "w": 0.8, "h": 0.8, "fill": "emerald"},
            {"type": "rect", "x": 1.9, "y": 5.4, "w": 1.6, "h": 0.4, "fill": "sky"}]},
    ]
    out = build_deck(_spec([_freeform_slide(prims)]), tmp_path / "freeform.pptx")
    assert out.exists()

    from pptx import Presentation
    prs = Presentation(str(out))
    slide = prs.slides[0]
    shape_types = [s.shape_type for s in slide.shapes]
    assert any(s.has_table for s in slide.shapes if hasattr(s, "has_table"))
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert MSO_SHAPE_TYPE.GROUP in shape_types


@pytest.mark.parametrize("chart_type,extra", [
    ("column", {}),
    ("stacked-bar", {}),
    ("line", {}),
    ("pie", {"slice_colors": ["emerald", "orange"]}),
    ("donut", {}),
    ("area", {}),
])
def test_build_chart_native_category_types(tmp_path, chart_type, extra):
    fill = {"chart_type": chart_type, "categories": ["a", "b"],
            "series": [{"name": "s", "values": [3, 5], "color": "emerald"}]}
    fill.update(extra)
    out = build_deck(_spec([_chart_slide(fill)]), tmp_path / f"{chart_type}.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert any(getattr(s, "has_chart", False) for s in prs.slides[0].shapes)


def test_build_chart_native_scatter(tmp_path):
    fill = {"chart_type": "scatter",
            "series": [{"name": "s", "x_values": [1, 2, 3], "y_values": [2, 4, 8], "color": "royal"}]}
    out = build_deck(_spec([_chart_slide(fill)]), tmp_path / "scatter.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert any(getattr(s, "has_chart", False) for s in prs.slides[0].shapes)


def test_chart_native_series_fill_uses_schemeclr(tmp_path):
    fill = {"chart_type": "column", "categories": ["a", "b"],
            "series": [{"name": "s", "values": [3, 5], "color": "emerald"}]}
    out = build_deck(_spec([_chart_slide(fill)]), tmp_path / "scheme.pptx")
    from pptx import Presentation
    prs = Presentation(str(out))
    chart = next(s for s in prs.slides[0].shapes if getattr(s, "has_chart", False)).chart
    xml = chart.series[0]._element.xml
    assert "schemeClr" in xml and "srgbClr" not in xml.split("schemeClr")[0][-200:]
