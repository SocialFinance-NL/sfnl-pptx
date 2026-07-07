"""Build-layer integration: fixture deck -> valid .pptx via node build_deck.js."""
import shutil
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
ENGINE = ROOT / "engine"
BUILD = ENGINE / "web" / "build"
FIXTURES = Path(__file__).parent / "fixtures"
NODE = shutil.which("node")
SEMANTIC_FRAME_COLORS = ("F87F4F", "3B62C1", "45B6E2", "6AC6BA", "F95D63")

pytestmark = pytest.mark.skipif(
    NODE is None or not (BUILD / "node_modules").exists(),
    reason="node or web build deps unavailable",
)


def _workspace(dst_dir, name):
    ws = dst_dir / name
    shutil.copytree(FIXTURES / name, ws)
    shutil.copy(ENGINE / "web" / "sfnl.css", ws / "slides" / "sfnl.css")
    shutil.copytree(ENGINE / "web" / "assets" / "chrome", ws / "slides" / "chrome")
    return ws


def _build(ws):
    return subprocess.run([NODE, str(BUILD / "build_deck.js"), str(ws)],
                          capture_output=True, text=True, timeout=300)


@pytest.fixture(scope="module")
def built(tmp_path_factory):
    ws = _workspace(tmp_path_factory.mktemp("webdeck"), "webdeck")
    res = _build(ws)
    assert res.returncode == 0, res.stderr + res.stdout
    out = ws / "webdeck.pptx"
    assert out.exists()
    return Presentation(str(out))


def test_all_slides_present(built):
    assert len(list(built.slides)) == 3


def test_build_embeds_official_layout_gallery(built):
    assert len(built.slide_masters) == 3
    assert sum(len(master.slide_layouts) for master in built.slide_masters) == 31


def test_title_chrome_lands_on_content_slide(built):
    slide = list(built.slides)[1]
    titles = [s for s in slide.shapes if s.has_text_frame and "DRIE CIJFERS" in s.text_frame.text]
    assert titles, "title text missing"
    assert titles[0].top < Inches(0.9)
    allowed = {"Lato Light", "Montserrat Light"}
    fonts = {r.font.name for p in titles[0].text_frame.paragraphs for r in p.runs if r.text.strip()}
    assert fonts, "title has no visible runs"
    assert fonts <= allowed
    assert titles[0].text_frame.text == titles[0].text_frame.text.upper()


def test_content_slide_has_brand_frame_marker(built):
    slide = list(built.slides)[1]
    frame_markers = [
        shape
        for shape in slide.shapes
        if (not shape.has_text_frame or not shape.text_frame.text.strip())
        and (shape.width >= Inches(1.0) or shape.height >= Inches(0.5))
        and any(color in shape._element.xml for color in SEMANTIC_FRAME_COLORS)
    ]
    assert frame_markers, "structural brand frame marker missing"


def test_notes_carry_dossier_refs(built):
    slide = list(built.slides)[1]
    assert slide.has_notes_slide
    assert "R1" in slide.notes_slide.notes_text_frame.text


def test_chart_injected_with_brand_colors(built):
    slide = list(built.slides)[2]
    frames = [s for s in slide.shapes if s.has_chart]
    assert frames, "native chart missing"
    assert "F87F4F" in frames[0].chart._chartSpace.xml


def test_page_number_and_logo_chrome(built):
    slides = list(built.slides)
    # cover has chrome none -> no page number '1' textbox bottom-right
    assert not any(s.has_text_frame and s.text_frame.text.strip() == "1" for s in slides[0].shapes)
    two = [s for s in slides[1].shapes if s.has_text_frame and s.text_frame.text.strip() == "2"]
    assert two and two[0].top > Inches(5.0)
    assert any(s.shape_type == 13 for s in slides[1].shapes)  # PICTURE = logo


def test_editorial_fixture_builds(tmp_path):
    ws = _workspace(tmp_path, "webdeck-editorial")
    res = _build(ws)
    assert res.returncode == 0, res.stderr + res.stdout
    prs = Presentation(str(ws / "webdeck-editorial.pptx"))
    assert len(prs.slides) == 3
    assert any(s.has_chart for s in prs.slides[2].shapes)


def test_overflow_deck_fails_loudly(tmp_path):
    ws = _workspace(tmp_path, "webdeck-overflow")
    res = _build(ws)
    assert res.returncode != 0
    assert "overflows body" in (res.stderr + res.stdout)


@pytest.fixture(scope="module")
def built_visuals(tmp_path_factory):
    ws = _workspace(tmp_path_factory.mktemp("webdeck-visuals"), "webdeck-visuals")
    res = _build(ws)
    assert res.returncode == 0, res.stderr + res.stdout
    return Presentation(str(ws / "webdeck-visuals.pptx"))


def test_visuals_table_is_native(built_visuals):
    slide = list(built_visuals.slides)[0]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    tbl = tables[0].table
    assert len(tbl.rows) == 6
    assert "RESULTAAT" not in tbl.cell(0, 0).text  # titel staat buiten de tabel
    assert "679.745" in tbl.cell(2, 2).text


def test_visuals_diagram_has_shapes_and_connectors(built_visuals):
    import re
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    slide = list(built_visuals.slides)[1]
    autos = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]

    # pptxgenjs serialiseert connectorlijnen NIET als MSO_SHAPE_TYPE.LINE maar
    # als AUTO_SHAPE met <a:prstGeom prst="line"/>: python-pptx classificeert
    # elke <p:sp> als AUTO_SHAPE, los van de onderliggende prstGeom-waarde. Er
    # is geen enkele echte MSO_SHAPE_TYPE.LINE shape in de output; de
    # prstGeom-waarde is de betrouwbare marker voor "dit is een connectorlijn".
    assert not any(s.shape_type == MSO_SHAPE_TYPE.LINE for s in slide.shapes), (
        "verrassing: pptxgenjs levert nu wel MSO_SHAPE_TYPE.LINE op, "
        "asserties hieronder moeten herzien worden"
    )
    line_prst = re.compile(r'prstGeom\s+prst="line"')
    line_segments = [s for s in autos if line_prst.search(s._element.xml)]
    shapes_only = [s for s in autos if not line_prst.search(s._element.xml)]

    # chevron + 3 nodes, geteld ZONDER de connectorsegmenten (die komen ook
    # als AUTO_SHAPE binnen en zouden verdwenen nodes anders maskeren)
    assert len(shapes_only) >= 4, (
        f"expected >=4 non-line autoshapes (chevron + 3 nodes), "
        f"got {len(shapes_only)}: {[s.name for s in shapes_only]}"
    )
    assert len(line_segments) >= 4, (
        f"expected >=4 line-geometry autoshapes (straight(1) + elbow(3) "
        f"connector segments), got {len(line_segments)}: "
        f"{[s.name for s in line_segments]}"
    )
