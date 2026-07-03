from pptx import Presentation
from pptx.util import Inches

from scripts.rearrange import rearrange


def _deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    for title in ("A", "B", "C"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5)).text_frame.text = title
    return prs


def _titles(prs):
    return [slide.shapes[0].text_frame.text for slide in prs.slides]


def test_rearrange_reorders_duplicates_and_deletes_slides(tmp_path):
    prs = _deck()
    rearrange(prs, [2, 0, 2])
    out = tmp_path / "out.pptx"
    prs.save(out)

    reopened = Presentation(str(out))
    assert _titles(reopened) == ["C", "A", "C"]
